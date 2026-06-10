"""
Smart WSR Agent – Core Engine
All data processing, report building, export, LLM, and agent logic.
Shared across all pages via import.
"""

import os
import re
import io
import json
import math
import textwrap
import tempfile
import socket
import smtplib
import colorsys
from pathlib import Path
from datetime import datetime, timedelta

import numpy as np
import pandas as pd
import streamlit as st

# PDF
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

# PPT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Image
from PIL import Image, ImageDraw, ImageFont

# Email
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
from email.mime.image import MIMEImage

# Charts
import plotly.express as px
import plotly.graph_objects as go

# LLM
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.messages import HumanMessage

# ============================================================
# CONSTANTS
# ============================================================

APP_NAME = "Weekly Status Report"
BASE_DIR = Path.cwd()
OUT_DIR = BASE_DIR / "outputs"
OUT_DIR.mkdir(exist_ok=True)

CLR_GREEN  = colors.HexColor("#1a7f4b")   # BusinessNext green accent
CLR_PINK   = colors.HexColor("#eef2fb")   # light blue tint for alternating rows
CLR_BORDER = colors.HexColor("#e0e4ed")   # border grey
CLR_TOTAL  = colors.HexColor("#d0d6f0")   # total row – muted blue-grey
CLR_HEADER = colors.HexColor("#1a3a8c")   # BusinessNext royal blue (table headers)

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")

llm = ChatGoogleGenerativeAI(
    model="gemini-2.0-flash",
    temperature=0.7,
    api_key="AIzaSyDbPc3_rAqopnFDiA8-5On21aD6rbOLAvY"
)

# ============================================================
# SESSION STATE
# ============================================================

def init_state():
    defaults = {
        "source_df": None,
        "working_df": None,
        "report": None,
        "chat_history": [],
        "last_prompt": "",
        "awaiting_email": False,
        "email_prompt": "",
        "week_mode": "last",
        "smtp_host": "smtp.office365.com",
        "smtp_port": "587",
        "smtp_user": "",
        "smtp_pass": "",
        "smtp_sender": "umesh.tiwari@businessnext.com",
        "project_query": "",
        "project_details": "",
        "export_file": None,
        "export_filename": None,
        "show_download": False,
        "awaiting_next_action": False,
        "quick_project_search": "",
        # ── config-driven workflow ─────────────────────────────
        "col_config": {
            "project": None,
            "status": None,
            "dates": [],
            "display": [],
            "metrics": [],
        },
        "section_configs": [],   # list of section dicts
        "filter_rules": [],      # list of filter rule dicts
        "col_config_initialised": False,  # guard: auto-detect runs only once per data load
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v


def add_chat(role, msg):
    st.session_state.chat_history.append({"role": role, "msg": msg})


def reset_all():
    st.session_state.chat_history = []
    st.session_state.source_df = None
    st.session_state.working_df = None
    st.session_state.report = None
    st.session_state.awaiting_next_action = False
    st.session_state.awaiting_email = False
    st.session_state.project_query = ""
    st.session_state.project_details = ""
    st.session_state.export_file = None
    st.session_state.export_filename = None
    st.session_state.show_download = False
    st.session_state.col_config = {
        "project": None, "status": None, "dates": [], "display": [], "metrics": [],
    }
    st.session_state.section_configs = []
    st.session_state.filter_rules = []
    st.session_state.col_config_initialised = False

# ============================================================
# CSS STYLES
# ============================================================

def inject_styles():
    st.markdown(
        """
        <style>
        /* ── BusinessNext Enterprise Theme ─────────────────────── */

        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

        html, body, [class*="css"], .stApp {
            font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif !important;
        }

        /* Page background */
        .stApp { background-color: #f0f2f6 !important; }
        .stMain > div { padding-top: 1.5rem; }

        /* ── Sidebar ───────────────────────────────────────────── */
        section[data-testid="stSidebar"] {
            background-color: #0a0a1e !important;
            border-right: 1px solid #1a1a3a !important;
        }
        section[data-testid="stSidebar"] * { color: #c8cde8 !important; }
        section[data-testid="stSidebar"] h1,
        section[data-testid="stSidebar"] h2,
        section[data-testid="stSidebar"] h3 { color: #ffffff !important; }
        section[data-testid="stSidebar"] strong { color: #e8eaf6 !important; }

        /* ── App header ────────────────────────────────────────── */
        .app-header {
            background: linear-gradient(135deg, #0a0a1e 0%, #1a3a8c 100%);
            color: white;
            padding: 22px 28px;
            border-radius: 10px;
            margin-bottom: 22px;
            box-shadow: 0 4px 20px rgba(10, 10, 30, 0.28);
            display: flex;
            align-items: center;
        }
        .app-header h1 {
            margin: 0;
            font-size: 1.9rem;
            font-weight: 700;
            letter-spacing: 0.01em;
            color: #ffffff;
        }
        .app-header .subtitle {
            margin: 6px 0 0;
            color: rgba(255,255,255,0.72);
            font-size: 0.93rem;
        }

        /* ── Section cards ─────────────────────────────────────── */
        .section-card {
            background: #ffffff;
            border: 1px solid #e0e4ed;
            border-left: 4px solid #1a3a8c;
            border-radius: 8px;
            padding: 15px 20px;
            margin-bottom: 16px;
            box-shadow: 0 2px 8px rgba(10, 10, 30, 0.05);
        }
        .section-card h3 {
            margin: 0;
            color: #0a0a1e;
            font-size: 1.05rem;
            font-weight: 600;
        }
        .section-card .section-note {
            margin-top: 5px;
            color: #5a6275;
            font-size: 0.88rem;
        }

        /* ── Report card ───────────────────────────────────────── */
        .report-card {
            background: #ffffff;
            border: 1px solid #e0e4ed;
            border-top: 4px solid #1a3a8c;
            border-radius: 8px;
            padding: 20px 24px;
            margin-bottom: 22px;
            box-shadow: 0 2px 10px rgba(10, 10, 30, 0.07);
        }
        .report-card .report-label {
            font-size: 0.8rem;
            color: #6b7280;
            text-transform: uppercase;
            letter-spacing: 0.06em;
            margin-bottom: 6px;
        }
        .report-card h2 {
            margin: 0;
            color: #0a0a1e;
            font-size: 1.55rem;
            font-weight: 700;
        }
        .report-card .subtitle { color: #5a6275; font-size: 0.9rem; margin-top: 6px; }

        /* ── Tabs ──────────────────────────────────────────────── */
        .stTabs [data-baseweb="tab-list"] {
            gap: 8px;
            background: transparent !important;
            border-bottom: 2px solid #e0e4ed !important;
            padding-bottom: 0 !important;
        }
        .stTabs [data-baseweb="tab"] {
            border-radius: 6px 6px 0 0 !important;
            padding: 9px 22px !important;
            font-weight: 600 !important;
            font-size: 0.88rem !important;
            background: #f0f2f6 !important;
            color: #5a6275 !important;
            border: 1px solid #e0e4ed !important;
            border-bottom: none !important;
            margin-bottom: -2px !important;
        }
        .stTabs [aria-selected="true"][data-baseweb="tab"] {
            background: #0a0a1e !important;
            color: #ffffff !important;
            border-color: #0a0a1e !important;
        }
        .stTabs [data-baseweb="tab-highlight"] { display: none !important; }
        .stTabs [data-baseweb="tab-panel"] {
            background: #ffffff;
            border: 1px solid #e0e4ed;
            border-top: none;
            border-radius: 0 0 8px 8px;
            padding: 20px !important;
        }

        /* ── Buttons ───────────────────────────────────────────── */
        .stButton > button,
        .stDownloadButton > button {
            border-radius: 999px !important;
            font-weight: 600 !important;
            font-size: 0.88rem !important;
            padding: 0.55rem 1.4rem !important;
            transition: all 0.18s ease !important;
            letter-spacing: 0.01em !important;
        }
        button[data-testid="baseButton-primary"],
        .stButton > button[kind="primary"] {
            background: #0a0a1e !important;
            color: #ffffff !important;
            border: 2px solid #0a0a1e !important;
            box-shadow: 0 4px 12px rgba(10, 10, 30, 0.22) !important;
        }
        button[data-testid="baseButton-primary"]:hover,
        .stButton > button[kind="primary"]:hover {
            background: #1a3a8c !important;
            border-color: #1a3a8c !important;
            box-shadow: 0 6px 18px rgba(26, 58, 140, 0.32) !important;
        }
        button[data-testid="baseButton-secondary"],
        .stButton > button[kind="secondary"] {
            background: #ffffff !important;
            color: #0a0a1e !important;
            border: 2px solid #0a0a1e !important;
        }
        button[data-testid="baseButton-secondary"]:hover,
        .stButton > button[kind="secondary"]:hover {
            background: #f0f2f6 !important;
        }
        .stDownloadButton > button {
            background: #1a3a8c !important;
            color: #ffffff !important;
            border: 2px solid #1a3a8c !important;
        }
        .stDownloadButton > button:hover {
            background: #0a0a1e !important;
            border-color: #0a0a1e !important;
        }

        /* ── Link button ───────────────────────────────────────── */
        .stLinkButton a {
            border-radius: 999px !important;
            font-weight: 600 !important;
            font-size: 0.88rem !important;
        }
        [data-testid="stLinkButton"] a[kind="primary"] {
            background: #0a0a1e !important;
            color: #ffffff !important;
        }

        /* ── Inputs / selects ──────────────────────────────────── */
        .stTextInput > div > div > input,
        .stTextArea > div > textarea {
            border-radius: 7px !important;
            border: 1px solid #d1d9e6 !important;
            background: #ffffff !important;
            font-size: 0.9rem !important;
        }
        .stTextInput > div > div > input:focus,
        .stTextArea > div > textarea:focus {
            border-color: #1a3a8c !important;
            box-shadow: 0 0 0 2px rgba(26,58,140,0.15) !important;
        }
        .stSelectbox > div > div,
        .stMultiSelect > div > div {
            border-radius: 7px !important;
            border-color: #d1d9e6 !important;
        }

        /* ── DataFrames ────────────────────────────────────────── */
        .stDataFrame {
            border-radius: 8px !important;
            overflow: hidden !important;
            border: 1px solid #e0e4ed !important;
            box-shadow: 0 2px 8px rgba(10,10,30,0.05) !important;
        }
        /* Glide data grid header */
        .stDataFrame [data-testid="glideDataEditor"] .gdg-header,
        .stDataFrame canvas { border-radius: 0 !important; }

        /* ── Alerts / info boxes ───────────────────────────────── */
        .stAlert {
            border-radius: 7px !important;
            border-width: 1px !important;
        }
        [data-testid="stInfo"] {
            border-left: 4px solid #1a3a8c !important;
            background: #eef2fb !important;
            color: #0a0a1e !important;
        }
        [data-testid="stSuccess"] {
            border-left: 4px solid #1a7f4b !important;
        }
        [data-testid="stWarning"] {
            border-left: 4px solid #c07a0a !important;
        }
        [data-testid="stError"] {
            border-left: 4px solid #c0180a !important;
        }

        /* ── Divider ───────────────────────────────────────────── */
        hr { border-color: #e0e4ed !important; margin: 18px 0 !important; }

        /* ── Headings ──────────────────────────────────────────── */
        .stMarkdown h2, .stMarkdown h3, .stMarkdown h4 {
            color: #0a0a1e !important;
            font-weight: 700 !important;
        }

        /* ── Expander ──────────────────────────────────────────── */
        .stExpander {
            border: 1px solid #e0e4ed !important;
            border-radius: 8px !important;
            background: #ffffff !important;
        }
        .stExpander summary {
            font-weight: 600 !important;
            color: #0a0a1e !important;
        }

        /* ── Number inputs ─────────────────────────────────────── */
        .stNumberInput > div > div > input {
            border-radius: 7px !important;
            border-color: #d1d9e6 !important;
        }

        /* ════════════════════════════════════════════════════════
           SIDEBAR — Rich Enterprise Design
           ════════════════════════════════════════════════════════ */

        /* Logo / branding header */
        .sb-brand {
            background: linear-gradient(135deg, #0a0a1e 0%, #1a3a8c 100%);
            margin: -1rem -1rem 0 -1rem;
            padding: 22px 20px 18px;
            border-bottom: 1px solid rgba(255,255,255,0.08);
        }
        .sb-brand-title {
            font-size: 1.25rem;
            font-weight: 800;
            color: #ffffff !important;
            letter-spacing: 0.02em;
            display: flex;
            align-items: center;
            gap: 8px;
        }
        .sb-brand-sub {
            font-size: 0.76rem;
            color: rgba(255,255,255,0.55) !important;
            margin-top: 4px;
            letter-spacing: 0.04em;
            text-transform: uppercase;
        }
        .sb-brand-badge {
            display: inline-block;
            background: rgba(255,255,255,0.12);
            border: 1px solid rgba(255,255,255,0.18);
            border-radius: 4px;
            font-size: 0.65rem;
            color: rgba(255,255,255,0.7) !important;
            padding: 2px 7px;
            margin-top: 8px;
            letter-spacing: 0.06em;
            text-transform: uppercase;
        }

        /* Section label */
        .sb-label {
            font-size: 0.68rem;
            font-weight: 700;
            letter-spacing: 0.1em;
            text-transform: uppercase;
            color: #4a5380 !important;
            padding: 14px 0 6px;
        }

        /* Status cards */
        .sb-status {
            background: rgba(255,255,255,0.04);
            border: 1px solid rgba(255,255,255,0.08);
            border-radius: 8px;
            padding: 10px 13px;
            margin-bottom: 8px;
            display: flex;
            align-items: center;
            gap: 10px;
        }
        .sb-status-icon {
            font-size: 1.1rem;
            flex-shrink: 0;
        }
        .sb-status-text { flex: 1; }
        .sb-status-main {
            font-size: 0.85rem;
            font-weight: 600;
            color: #e8eaf6 !important;
        }
        .sb-status-sub {
            font-size: 0.75rem;
            color: #6b75a8 !important;
            margin-top: 1px;
        }
        .sb-status.success { border-color: rgba(26,127,75,0.35); background: rgba(26,127,75,0.08); }
        .sb-status.info    { border-color: rgba(26,58,140,0.35); background: rgba(26,58,140,0.10); }
        .sb-status.idle    { border-color: rgba(255,255,255,0.07); background: rgba(255,255,255,0.03); }

        /* Nav cards */
        .sb-nav-card {
            background: rgba(255,255,255,0.04);
            border: 1px solid rgba(255,255,255,0.07);
            border-radius: 8px;
            padding: 9px 13px;
            margin-bottom: 6px;
            cursor: default;
            transition: background 0.15s;
        }
        .sb-nav-card:hover { background: rgba(255,255,255,0.08); }
        .sb-nav-card-title {
            font-size: 0.87rem;
            font-weight: 600;
            color: #d0d6f0 !important;
        }
        .sb-nav-card-desc {
            font-size: 0.75rem;
            color: #5a6494 !important;
            margin-top: 2px;
            line-height: 1.4;
        }

        /* Workflow steps */
        .sb-step {
            display: flex;
            gap: 10px;
            padding: 8px 0;
            border-bottom: 1px solid rgba(255,255,255,0.05);
        }
        .sb-step:last-child { border-bottom: none; }
        .sb-step-num {
            width: 22px;
            height: 22px;
            border-radius: 50%;
            background: #1a3a8c;
            color: #ffffff !important;
            font-size: 0.72rem;
            font-weight: 700;
            display: flex;
            align-items: center;
            justify-content: center;
            flex-shrink: 0;
            margin-top: 1px;
        }
        .sb-step-body {}
        .sb-step-title {
            font-size: 0.82rem;
            font-weight: 600;
            color: #c8cde8 !important;
        }
        .sb-step-desc {
            font-size: 0.75rem;
            color: #5a6494 !important;
            line-height: 1.45;
            margin-top: 2px;
        }

        /* Help tip */
        .sb-tip {
            background: rgba(26,58,140,0.15);
            border: 1px solid rgba(26,58,140,0.3);
            border-radius: 7px;
            padding: 9px 12px;
            font-size: 0.76rem;
            color: #9aa8d8 !important;
            line-height: 1.5;
            margin-top: 6px;
        }
        .sb-tip strong { color: #c0cbf0 !important; }

        /* Sidebar divider */
        .sb-divider {
            border: none;
            border-top: 1px solid rgba(255,255,255,0.07);
            margin: 12px 0;
        }

        /* Sidebar expander override */
        section[data-testid="stSidebar"] .stExpander {
            background: rgba(255,255,255,0.03) !important;
            border: 1px solid rgba(255,255,255,0.08) !important;
            border-radius: 8px !important;
            margin-bottom: 6px !important;
        }
        section[data-testid="stSidebar"] .stExpander summary {
            color: #c8cde8 !important;
            font-size: 0.83rem !important;
            font-weight: 600 !important;
        }
        section[data-testid="stSidebar"] .stExpander summary:hover {
            color: #ffffff !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

# ============================================================
# DATE HELPERS
# ============================================================

def today():
    return pd.Timestamp.today().normalize()


def week_range(mode="last"):
    d = today()
    m = str(mode).lower()
    if "current" in m or "this" in m:
        start = d - timedelta(days=d.weekday())
    else:
        start = d - timedelta(days=d.weekday() + 7)
    end = start + timedelta(days=4)
    return f"{start:%d %b} - {end:%d %b, %Y}"


def week_window(mode="last"):
    d = today()
    m = str(mode).lower()
    if "current" in m or "this" in m:
        start = d - timedelta(days=d.weekday())
    else:
        start = d - timedelta(days=d.weekday() + 7)
    end = start + timedelta(days=4)
    return start, end


def previous_week_range():
    return week_range("last")


def get_report_week_range(report):
    if report is None:
        return previous_week_range()
    return report.get("week_range", previous_week_range())


def get_report_title(report):
    return f"Weekly Status Report : {get_report_week_range(report)}"


def output_folder():
    f = OUT_DIR / datetime.now().strftime("%Y%m%d_%H%M%S")
    f.mkdir(exist_ok=True)
    return f

# ============================================================
# DATA HELPERS
# ============================================================

def to_dt(s):
    return pd.to_datetime(s, errors="coerce", dayfirst=True)


def clean_text(x):
    return str(x).strip()


def normalize_col_key(col):
    key = str(col).strip().lower()
    key = re.sub(r"[\s_\-]+", " ", key)
    key = re.sub(r"[^a-z0-9 ]", "", key)
    return key


def is_string_column(series):
    return (
        pd.api.types.is_string_dtype(series)
        or pd.api.types.is_object_dtype(series)
        or pd.api.types.is_categorical_dtype(series)
    )


def find_col(df, names):
    if df is None or len(df) == 0:
        return None
    if isinstance(names, str):
        names = [names]
    df_cols = [str(c).strip() for c in df.columns]
    norm_cols = {normalize_col_key(c): c for c in df_cols}
    for n in names:
        n_norm = normalize_col_key(n)
        if n_norm in norm_cols:
            return norm_cols[n_norm]
    for n in names:
        n_norm = normalize_col_key(n)
        n_tokens = set(n_norm.split())
        for c_norm, c in norm_cols.items():
            if n_norm in c_norm or c_norm in n_norm:
                return c
            if n_tokens and n_tokens.issubset(set(c_norm.split())):
                return c
            if len(n_tokens) == 1 and c_norm.startswith(next(iter(n_tokens))):
                return c
    return None


def guess_status_column(df):
    if df is None or len(df) == 0:
        return None
    keywords = [
        "uat", "signed off", "production", "go live", "golive",
        "development", "dev", "ud", "hold", "cutover", "stage", "state",
        "progress", "in progress", "current", "planned", "pending", "blocked",
        "complete", "completed", "testing", "deployment", "release"
    ]
    status_headers = [
        c for c in df.columns
        if any(k in str(c).lower() for k in ["status", "current", "stage", "state", "phase"])
    ]
    if status_headers:
        best_header, best_score = None, -1
        for c in status_headers:
            if not is_string_column(df[c]):
                continue
            values = df[c].astype(str).str.strip().replace("nan", "", regex=False)
            non_empty = values[values != ""]
            score = len(non_empty.unique()) if len(non_empty) > 0 else 0
            if score > best_score:
                best_score = score
                best_header = c
        if best_header is not None:
            return best_header
    col = find_col(df, ["Status", "Status Code", "Current Status", "Project Status", "Status Cac", "Stage", "State"])
    if col:
        return col
    best_col, best_score = None, -1
    for c in df.columns:
        if is_string_column(df[c]):
            values = df[c].astype(str).str.lower().fillna("")
            score = sum(values.str.contains(kw, na=False).sum() for kw in keywords)
            if score > best_score:
                best_score = score
                best_col = c
    if best_score > 0:
        return best_col
    candidate, best_ratio = None, 0
    for c in df.columns:
        if is_string_column(df[c]):
            values = df[c].astype(str).str.strip().replace("nan", "", regex=False)
            non_empty = values[values != ""]
            if len(non_empty) == 0:
                continue
            unique_count = non_empty.nunique()
            ratio = unique_count / len(non_empty)
            if unique_count > 1 and ratio > best_ratio:
                best_ratio = ratio
                candidate = c
    return candidate


def guess_project_column(df):
    if df is None or len(df) == 0:
        return None
    project = find_col(df, [
        "Project", "Project Name", "Project Na", "Project Module",
        "Module", "Module Name", "Project/Module", "Initiative", "Program"
    ])
    if project:
        return project
    project_headers = [
        c for c in df.columns
        if any(k in str(c).lower() for k in ["project", "module", "initiative", "program", "workstream"])
    ]
    if project_headers:
        best_header, best_score = None, -1
        for c in project_headers:
            if not is_string_column(df[c]):
                continue
            values = df[c].astype(str).str.strip().replace("nan", "", regex=False)
            non_empty = values[values != ""]
            score = len(non_empty.unique()) if len(non_empty) > 0 else 0
            if score > best_score:
                best_score = score
                best_header = c
        if best_header is not None:
            return best_header
    candidate, best_ratio = None, 0
    for c in df.columns:
        if is_string_column(df[c]):
            values = df[c].astype(str).str.strip().replace("nan", "", regex=False)
            non_empty = values[values != ""]
            if len(non_empty) == 0:
                continue
            unique_count = non_empty.nunique()
            ratio = unique_count / len(non_empty)
            if unique_count > 1 and ratio > best_ratio:
                best_ratio = ratio
                candidate = c
    return candidate


def guess_date_column(df, keywords):
    if df is None or len(df) == 0:
        return None
    col = find_col(df, keywords)
    if col:
        return col
    norm_keywords = [normalize_col_key(k) for k in keywords]
    for c in df.columns:
        c_norm = normalize_col_key(c)
        if any(k in c_norm for k in norm_keywords):
            return c
    preferred = ["uat", "sit", "go live", "golive", "production", "actual", "created", "planned", "date"]
    header_candidates = []
    for c in df.columns:
        c_norm = normalize_col_key(c)
        score = sum(1 for tk in preferred if tk in c_norm)
        if score > 0:
            header_candidates.append((score, c))
    if header_candidates:
        header_candidates.sort(key=lambda x: (-x[0], str(x[1])))
        return header_candidates[0][1]
    best_col, best_score = None, -1
    row_count = len(df)
    for c in df.columns:
        if is_string_column(df[c]) or pd.api.types.is_datetime64_any_dtype(df[c]):
            parsed = pd.to_datetime(df[c], errors='coerce', dayfirst=True)
            score = parsed.notna().sum()
            if score > best_score:
                best_score = score
                best_col = c
    if best_score >= max(1, row_count // 8):
        return best_col
    if best_score > 0:
        return best_col
    for c in df.columns:
        if pd.api.types.is_datetime64_any_dtype(df[c]):
            return c
    return None


def std_cols(df):
    return {
        "project": find_col(df, [
            "Project Module", "Project/Module", "Project Name", "Project Na",
            "Project", "Module", "Module Name", "Project Name "
        ]),
        "status": find_col(df, [
            "Status Code", "Status", "Status Cac", "Current Status", "Project Status", "Status "
        ]),
        "uat": find_col(df, [
            "Actual UAT Start", "UAT Start", "UAT Date", "UAT Start Date", "UAT",
            "Actual SIT Date", "Actual SIT End Date", "Actual SIT PreUAT Date",
            "Actual SIT PreUAT De", "Actual SIT"
        ]),
        "golive": find_col(df, [
            "Actual Go live", "Actual Go Live", "Go Live Date", "Go Live",
            "Production Date", "Go-Live Date", "Actual Go Created",
            "Actual Go Created On", "Actual Go Created O", "Actual GO", "Go Created"
        ]),
    }


def load_file(file):
    file.seek(0)
    name = file.name.lower()

    if name.endswith(".csv"):
        for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
            try:
                file.seek(0)
                df = pd.read_csv(file, encoding=enc)
                if not df.empty:
                    return df
            except Exception:
                continue
        raise ValueError("Could not read CSV — file may be empty or use an unsupported encoding.")

    if name.endswith(".xlsx") or name.endswith(".xls"):
        file.seek(0)
        sheets = pd.read_excel(file, sheet_name=None)
        for _, df in sheets.items():
            df = df.dropna(how="all").reset_index(drop=True)
            if len(df) > 0 and len(df.columns) > 0:
                return df
        raise ValueError("Excel file has no data in any sheet.")

    raise ValueError(f"Unsupported file type: {file.name}")


def safe_str(x):
    if pd.isna(x):
        return ""
    try:
        if isinstance(x, float) and x.is_integer():
            return str(int(x))
        return str(x)
    except:
        return ""


def fmt_date(val):
    try:
        d = to_dt(val)
        if pd.isna(d):
            return ""
        return d.strftime("%d-%b-%Y")
    except:
        return ""


def is_date_col(col):
    c = str(col).lower()
    return any(k in c for k in ["date", "start", "end", "live", "uat", "sit", "delivery", "created", "updated"])


def aging_bucket(age):
    try:
        age = int(float(age))
    except:
        age = 0
    if age <= 30:
        return "0-30 Days"
    elif age <= 60:
        return "30-60 Days"
    elif age <= 90:
        return "60-90 Days"
    return ">90 Days"


def norm(txt):
    return str(txt).lower().strip()

# ============================================================
# DATA PREPARATION
# ============================================================

def prepare_data(df):
    work = df.copy()
    cols = std_cols(work)
    status_col = cols.get("status") or guess_status_column(work)
    if not status_col:
        for c in work.columns:
            if "status" in str(c).lower() or "current" in str(c).lower():
                status_col = c
                break
    if not status_col:
        for c in work.columns:
            if work[c].dtype == 'object' or pd.api.types.is_string_dtype(work[c]):
                status_col = c
                break
    if not status_col:
        return work

    t = today()
    dev_col = find_col(work, ["Actual Dev Delivery Start"])
    uat_col = guess_date_column(work, [
        "Actual UAT Start", "UAT Start", "UAT Date", "UAT Start Date", "UAT",
        "Actual UAT Created O", "Actual UAT Created On", "Actual SIT End Date",
        "Actual SIT PreUAT Date"
    ])
    fsd_col = find_col(work, ["Actual FSD Delivery"])
    go_col = guess_date_column(work, [
        "Actual Go Live", "Actual Go live", "Go Live Date", "Go Live",
        "Production Date", "Actual Go Created On", "Actual Go Created O", "Actual Go Created"
    ])
    date_candidates = [c for c in work.columns if is_date_col(c)]

    aging_dates, aging_days = [], []

    def fallback_row_dt(row):
        for c in date_candidates:
            try:
                dt = pd.to_datetime(row[c], errors="coerce", dayfirst=True)
                if pd.notna(dt):
                    return dt
            except:
                continue
        return pd.NaT

    for _, r in work.iterrows():
        s = str(r[status_col]).lower().strip()
        dt = pd.NaT
        if any(x in s for x in ["development", "dev", "build", "coding", "in progress"]):
            if dev_col:
                dt = pd.to_datetime(r[dev_col], errors="coerce", dayfirst=True)
        elif any(x in s for x in ["uat signed off", "uat", "sit", "testing", "validation"]):
            if uat_col:
                dt = pd.to_datetime(r[uat_col], errors="coerce", dayfirst=True)
            elif fsd_col:
                dt = pd.to_datetime(r[fsd_col], errors="coerce", dayfirst=True)
        elif any(x in s for x in ["ud submitted", "ud signed off", "in ud", "user delivery", "ud"]):
            if fsd_col:
                dt = pd.to_datetime(r[fsd_col], errors="coerce", dayfirst=True)
        elif any(x in s for x in ["production cutover", "go live", "golive", "prod", "deployment", "release"]):
            if go_col:
                dt = pd.to_datetime(r[go_col], errors="coerce", dayfirst=True)
        elif any(x in s for x in ["pre-engagement", "pre engagement", "not started", "planning"]):
            if fsd_col:
                dt = pd.to_datetime(r[fsd_col], errors="coerce", dayfirst=True)
        if pd.isna(dt):
            dt = fallback_row_dt(r)
        if pd.isna(dt):
            aging_dates.append("")
            aging_days.append(0)
        else:
            dt = dt.normalize()
            diff = max(0, (t - dt).days)
            aging_dates.append(fmt_date(dt))
            aging_days.append(diff)

    work["Aging From Date"] = aging_dates
    work["Aging"] = aging_days
    for c in work.columns:
        if is_date_col(c):
            work[c] = work[c].apply(fmt_date)
    return work


def add_aging(df):
    work = df.copy()
    cols = std_cols(work)
    status = cols.get("status")
    uat = cols.get("uat")
    golive = cols.get("golive")
    if not status:
        for c in work.columns:
            if "status" in str(c).lower() or "current" in str(c).lower():
                status = c
                break
    if not status:
        for c in work.columns:
            if work[c].dtype == 'object':
                status = c
                break
    if not status:
        work["Status"] = ""
        status = "Status"
    s = work[status].astype(str).str.lower()
    work["Aging From Date"] = pd.NaT
    if uat:
        mask = s.str.contains("uat", na=False)
        work.loc[mask, "Aging From Date"] = to_dt(work.loc[mask, uat])
    if golive:
        mask = s.str.contains("production cutover", na=False)
        work.loc[mask, "Aging From Date"] = to_dt(work.loc[mask, golive])
    if uat:
        work["Aging From Date"] = work["Aging From Date"].fillna(to_dt(work[uat]))
    work["Aging"] = (today() - work["Aging From Date"]).dt.days
    work["Aging"] = (
        pd.to_numeric(work["Aging"], errors="coerce")
        .fillna(0).clip(lower=0).astype(int)
    )
    return work

# ============================================================
# REPORT BUILDERS
# ============================================================

def build_summary(work):
    cols = std_cols(work)
    status = cols.get("status") or guess_status_column(work)
    if not status:
        return pd.DataFrame()
    s = work[status].astype(str).fillna("").str.lower()
    data = {
        "Go Live": s.str.contains("production cutover|production|go live|golive|prod|deployment|release", na=False).sum(),
        "Dev": s.str.contains("development|dev|in progress|build|coding|design|planning", na=False).sum(),
        "PO Awaited": s.str.contains("ud signed off|po awaited|pending approval|approval awaited", na=False).sum(),
        "UAT": s.str.contains("uat|testing|test|validation|sit", na=False).sum(),
        "In UD": s.str.contains("ud submitted|in ud|user delivery|ud in progress|ud", na=False).sum(),
        "UD Not Started": s.str.contains("pre - engagement|pre engagement|preengagement|not started|planning", na=False).sum(),
        "On Hold": s.str.contains("hold|blocked|delay|delayed", na=False).sum(),
        "Grand Total": len(work),
    }
    return pd.DataFrame([data])


def build_aging(work):
    cols = std_cols(work)
    status = cols.get("status") or guess_status_column(work)
    if not status:
        return pd.DataFrame()
    temp = work.copy()
    temp[status] = temp[status].astype(str).fillna("").str.lower()
    pre_df = temp[temp[status].str.contains("pre - engagement", na=False)]
    main_df = temp[~temp[status].str.contains("pre - engagement", na=False)]

    def stage(x):
        if any(k in x for k in ["development", "dev", "build", "coding", "in progress"]):
            return "Dev"
        elif any(k in x for k in ["ud signed off", "po awaited", "pending approval", "approval awaited"]):
            return "PO Awaited"
        elif any(k in x for k in ["uat", "sit", "testing", "validation"]):
            return "UAT"
        elif any(k in x for k in ["ud submitted", "in ud", "user delivery", "ud"]):
            return "In UD"
        elif any(k in x for k in ["hold", "blocked", "delay", "delayed"]):
            return "On Hold"
        return "Other"

    def bucket(n):
        try:
            n = int(float(n))
        except:
            n = 0
        if n <= 30:
            return "0 - 30 Days"
        elif n <= 60:
            return "30 - 60 Days"
        elif n <= 90:
            return "60 - 90 Days"
        return "> 90 Days"

    main_df = main_df.copy()
    main_df["Stage"] = main_df[status].apply(stage)
    main_df["Bucket"] = main_df["Aging"].apply(bucket)
    ct = pd.crosstab(main_df["Stage"], main_df["Bucket"])
    rows_order = ["Dev", "PO Awaited", "UAT", "In UD", "On Hold"]
    cols_req = ["0 - 30 Days", "30 - 60 Days", "60 - 90 Days", "> 90 Days"]
    ct = ct.reindex(index=rows_order, columns=cols_req, fill_value=0)
    ct["UD Not Started"] = 0
    ct.loc["In UD", "UD Not Started"] = len(pre_df)
    ct = ct[["0 - 30 Days", "30 - 60 Days", "60 - 90 Days", "> 90 Days", "UD Not Started"]]
    ct["Grand Total"] = ct.sum(axis=1)
    ct.loc["Grand Total"] = ct.sum()
    return ct


def build_delivery(work, week_mode="last"):
    cols = std_cols(work)
    proj = cols.get("project") or guess_project_column(work)
    status = cols.get("status") or guess_status_column(work)
    uat_col = cols.get("uat") or guess_date_column(work, [
        "Actual UAT Start", "UAT Start", "UAT Date", "UAT Start Date", "UAT",
        "Actual UAT Created O", "Actual UAT Created On", "Actual SIT End Date", "Actual SIT PreUAT Date"
    ])
    go_col = cols.get("golive") or guess_date_column(work, [
        "Actual Go live", "Actual Go Live", "Go Live Date", "Go Live",
        "Production Date", "Actual Go Created On", "Actual Go Created O", "Actual Go Created"
    ])
    if not proj or not status:
        return pd.DataFrame()
    start, end = week_window(week_mode)
    out = []
    for _, r in work.iterrows():
        s = str(r[status]).lower()
        if uat_col and any(x in s for x in ["uat", "sit", "testing", "validation"]):
            dt = pd.to_datetime(r[uat_col], errors="coerce", dayfirst=True)
            if pd.notna(dt) and start <= dt.normalize() <= end:
                out.append({"Project Module": r[proj], "UAT / Go Live Date": dt.date(), "Status Code": r[status]})
        if go_col and any(x in s for x in ["production cutover", "go live", "golive", "prod", "deployment", "release"]):
            dt = pd.to_datetime(r[go_col], errors="coerce", dayfirst=True)
            if pd.notna(dt) and start <= dt.normalize() <= end:
                out.append({"Project Module": r[proj], "UAT / Go Live Date": dt.date(), "Status Code": r[status]})
    return pd.DataFrame(out)


def build_outlook(work):
    cols = std_cols(work)
    proj = cols.get("project") or guess_project_column(work)
    status = cols.get("status") or guess_status_column(work)
    uat_col = cols.get("uat") or guess_date_column(work, [
        "Actual UAT Start", "UAT Start", "UAT Date", "UAT Start Date", "UAT",
        "Actual UAT Created O", "Actual UAT Created On", "Actual SIT End Date", "Actual SIT PreUAT Date"
    ])
    go_col = find_col(work, ["Planned Go live", "Planned Go Live", "Planned Golive"]) or \
             cols.get("golive") or guess_date_column(work, [
        "Actual Go live", "Actual Go Live", "Go Live Date", "Go Live",
        "Production Date", "Actual Go Created On", "Actual Go Created O", "Actual Go Created"
    ])
    if not proj or not status:
        return pd.DataFrame()
    d = today()
    days_to_monday = (7 - d.weekday()) % 7
    if days_to_monday == 0:
        days_to_monday = 7
    next_mon = d + timedelta(days=days_to_monday)
    next_fri = next_mon + timedelta(days=4)
    out = []
    for _, r in work.iterrows():
        s = str(r[status]).strip().lower()
        if uat_col and any(x in s for x in ["development", "dev", "build", "coding", "in progress"]):
            dt = pd.to_datetime(r[uat_col], errors="coerce", dayfirst=True)
            if pd.notna(dt) and next_mon <= dt.normalize() <= next_fri:
                out.append({"Planned Delivery": "UAT Delivery", "Project Module": r[proj], "UAT or Go Live Date": dt.date()})
        if go_col and any(x in s for x in ["uat signedoff", "uat signed off", "production cutover", "go live", "golive", "prod", "deployment", "release"]):
            dt = pd.to_datetime(r[go_col], errors="coerce", dayfirst=True)
            if pd.notna(dt) and next_mon <= dt.normalize() <= next_fri:
                out.append({"Planned Delivery": "Go Live", "Project Module": r[proj], "UAT or Go Live Date": dt.date()})
    return pd.DataFrame(out)


def build_risks(work):
    cols = std_cols(work)
    proj = cols.get("project") or guess_project_column(work)
    status = cols.get("status") or guess_status_column(work)
    if not proj or not status:
        return pd.DataFrame()
    allow = ["ud submitted", "ud signed off", "uat", "uat signed off"]
    temp = work[work[status].astype(str).str.lower().apply(lambda x: any(a in x for a in allow))].copy()
    temp = temp.sort_values("Aging", ascending=False).head(3)
    return pd.DataFrame({
        "Project Module": temp[proj],
        "Status Code": temp[status],
        "Aging From Date": temp["Aging From Date"],
        "Aging (Days)": temp["Aging"],
    }).reset_index(drop=True)


def build_exec_summary(work, report):
    cols = std_cols(work)
    proj = cols.get("project") or guess_project_column(work)
    desc = find_col(work, ["Description", "Comments", "Latest Comment", "Remarks"])
    dt = find_col(work, ["Latest Date", "Updated Date", "Modified Date", "Comment Date", "Last Updated"])
    status = cols.get("status") or guess_status_column(work)
    if not proj or not desc:
        return ["No executive updates available."]
    temp = work.copy()
    if dt:
        temp["SortDate"] = pd.to_datetime(temp[dt], errors="coerce", dayfirst=True)
        temp = temp.sort_values("SortDate", ascending=False)
    if "Aging" in temp.columns:
        temp = temp.sort_values(["Aging", "SortDate"] if dt else ["Aging"], ascending=False)
    temp = temp.fillna("")
    bullets, used = [], set()
    for _, r in temp.iterrows():
        p = str(r[proj]).strip()
        if not p or p in used:
            continue
        used.add(p)
        raw = str(r[desc]).strip()
        stt = str(r[status]).strip() if status else ""
        aging = int(r["Aging"]) if "Aging" in r else 0
        txt = raw.lower()
        if any(x in txt for x in ["uat", "testing", "test case", "validation"]):
            msg = f"{p} is progressing in testing/UAT stage. Current focus is completion of validations and closure of pending observations."
        elif any(x in txt for x in ["go live", "deployment", "production"]):
            msg = f"{p} is nearing production readiness with deployment related activities in progress."
        elif any(x in txt for x in ["delay", "issue", "risk", "blocked", "pending"]):
            msg = f"{p} requires management attention due to pending dependencies/issues impacting timelines."
        elif any(x in txt for x in ["development", "build", "coding", "fix"]):
            msg = f"{p} is in active development phase with planned deliverables under execution."
        else:
            msg = f"{p} is currently in {stt} stage with ongoing progress updates under review."
        if aging >= 90:
            msg += f" Project aging is high ({aging} days) and needs priority tracking."
        bullets.append(msg)
        if len(bullets) >= 5:
            break
    if not bullets:
        bullets = ["Overall portfolio is progressing as planned with key deliveries under active monitoring."]
    return bullets


def build_report(df, week_mode="last"):
    work = prepare_data(df)
    report = {}
    report["working"] = work
    report["summary"] = build_summary(work)
    report["aging"] = build_aging(work)
    report["delivery"] = build_delivery(work, week_mode=week_mode)
    report["outlook"] = build_outlook(work)
    report["risks"] = build_risks(work)
    report["bullets"] = build_exec_summary(work, report)
    report["week_mode"] = week_mode
    report["week_range"] = week_range(week_mode)
    report["filter_criteria"] = report.get("filter_criteria", {})
    report["sections"] = None   # legacy mode — dashboard checks this
    return report


# ============================================================
# CONFIG-DRIVEN WORKFLOW
# ============================================================

def detect_col_config(df):
    """Auto-detect sensible column roles from a dataframe."""
    all_cols = list(df.columns)
    numeric_cols = [c for c in all_cols if pd.api.types.is_numeric_dtype(df[c])]
    date_cols = []
    for c in all_cols:
        if pd.api.types.is_datetime64_any_dtype(df[c]):
            date_cols.append(c)
        elif df[c].dtype == object:
            parsed = pd.to_datetime(df[c], errors="coerce", dayfirst=True)
            if parsed.notna().sum() >= max(1, len(df) // 4):
                date_cols.append(c)

    proj_col = guess_project_column(df)
    stat_col = guess_status_column(df)
    display = [c for c in all_cols if c not in numeric_cols][:10]

    return {
        "project": proj_col,
        "status": stat_col,
        "dates": date_cols[:4],
        "display": display,
        "metrics": numeric_cols[:6],
    }


def default_section_configs(df, col_config):
    """Build a sensible default section list from detected columns."""
    sections = []
    sid = 1

    status_col = col_config.get("status")
    project_col = col_config.get("project")
    display_cols = col_config.get("display") or list(df.columns)[:8]
    metric_cols = col_config.get("metrics") or []

    if status_col:
        sections.append({
            "id": f"s{sid}", "title": "Status Summary", "enabled": True,
            "section_type": "pivot", "group_by": status_col,
            "agg_func": "count", "agg_col": None,
            "display_cols": [], "sort_by": None, "sort_asc": False, "limit": None,
            "layout": "vertical",
        })
        sid += 1
        # Bar + Pie default side-by-side (horizontal)
        sections.append({
            "id": f"s{sid}", "title": "Status Distribution (Bar)", "enabled": True,
            "section_type": "bar_chart", "group_by": status_col,
            "agg_func": "count", "agg_col": None,
            "display_cols": [], "sort_by": None, "sort_asc": False, "limit": None,
            "layout": "horizontal",
        })
        sid += 1
        sections.append({
            "id": f"s{sid}", "title": "Status Share (Pie)", "enabled": True,
            "section_type": "pie_chart", "group_by": status_col,
            "agg_func": "count", "agg_col": None,
            "display_cols": [], "sort_by": None, "sort_asc": False, "limit": None,
            "layout": "horizontal",
        })
        sid += 1

    sections.append({
        "id": f"s{sid}", "title": "Key Metrics", "enabled": True,
        "section_type": "kpi_cards", "group_by": status_col,
        "agg_func": "count", "agg_col": None,
        "display_cols": metric_cols[:4] or display_cols[:4],
        "sort_by": None, "sort_asc": False, "limit": None,
        "layout": "vertical",
    })
    sid += 1

    sections.append({
        "id": f"s{sid}", "title": "Full Data Table", "enabled": True,
        "section_type": "table", "group_by": None,
        "agg_func": "count", "agg_col": None,
        "display_cols": display_cols, "sort_by": None, "sort_asc": True, "limit": 100,
        "layout": "vertical",
    })

    return sections


def apply_filter_rules(df, rules):
    """Apply a list of user-defined filter rules to df."""
    if not rules:
        return df
    masks = []
    connectors = []
    for r in rules:
        col = r.get("col")
        op = r.get("op", "contains")
        val = str(r.get("val", ""))
        if not col or col not in df.columns:
            continue
        try:
            if op == "equals":
                mask = df[col].astype(str).str.strip().str.lower() == val.strip().lower()
            elif op == "not_equals":
                mask = df[col].astype(str).str.strip().str.lower() != val.strip().lower()
            elif op == "contains":
                mask = df[col].astype(str).str.contains(val, case=False, na=False)
            elif op == "not_contains":
                mask = ~df[col].astype(str).str.contains(val, case=False, na=False)
            elif op == "greater_than":
                mask = pd.to_numeric(df[col], errors="coerce") > float(val)
            elif op == "less_than":
                mask = pd.to_numeric(df[col], errors="coerce") < float(val)
            elif op == "gte":
                mask = pd.to_numeric(df[col], errors="coerce") >= float(val)
            elif op == "lte":
                mask = pd.to_numeric(df[col], errors="coerce") <= float(val)
            elif op == "in_list":
                vals = [v.strip().lower() for v in val.split(",") if v.strip()]
                mask = df[col].astype(str).str.strip().str.lower().isin(vals)
            elif op == "not_in_list":
                vals = [v.strip().lower() for v in val.split(",") if v.strip()]
                mask = ~df[col].astype(str).str.strip().str.lower().isin(vals)
            elif op == "is_empty":
                mask = df[col].isna() | (df[col].astype(str).str.strip() == "")
            elif op == "is_not_empty":
                mask = ~(df[col].isna() | (df[col].astype(str).str.strip() == ""))
            elif op == "starts_with":
                mask = df[col].astype(str).str.lower().str.startswith(val.lower())
            elif op == "ends_with":
                mask = df[col].astype(str).str.lower().str.endswith(val.lower())
            else:
                continue
        except Exception:
            continue
        masks.append(mask)
        connectors.append(r.get("connector", "AND"))

    if not masks:
        return df
    result = masks[0]
    for i, m in enumerate(masks[1:], 1):
        if connectors[i] == "OR":
            result = result | m
        else:
            result = result & m
    return df[result].reset_index(drop=True)


def build_section_data(df, cfg):
    """Compute the data payload for a single section config."""
    stype = cfg.get("section_type", "table")
    group_by = cfg.get("group_by")
    agg_func = cfg.get("agg_func", "count")
    agg_col = cfg.get("agg_col")
    display_cols = [c for c in (cfg.get("display_cols") or []) if c in df.columns]
    sort_by = cfg.get("sort_by")
    sort_asc = cfg.get("sort_asc", False)
    limit = cfg.get("limit") or None

    try:
        if stype == "table":
            cols = display_cols if display_cols else list(df.columns)
            result = df[cols].copy()
            if sort_by and sort_by in result.columns:
                result = result.sort_values(sort_by, ascending=sort_asc)
            return result.head(limit) if limit else result

        if stype in ("pivot", "bar_chart", "pie_chart"):
            if not group_by or group_by not in df.columns:
                return pd.DataFrame({"Note": ["Select a Group By column"]})
            if agg_func == "count":
                result = df.groupby(group_by, dropna=False).size().reset_index(name="Count")
            elif agg_func == "nunique" and agg_col and agg_col in df.columns:
                result = df.groupby(group_by, dropna=False)[agg_col].nunique().reset_index(name="Distinct Count")
            elif agg_col and agg_col in df.columns:
                fn = {"sum": "sum", "mean": "mean", "min": "min", "max": "max"}.get(agg_func, "sum")
                result = df.groupby(group_by, dropna=False)[agg_col].agg(fn).reset_index()
                result.columns = [group_by, f"{fn.title()} of {agg_col}"]
            else:
                result = df.groupby(group_by, dropna=False).size().reset_index(name="Count")
            val_col = result.columns[1]
            result = result.sort_values(val_col, ascending=sort_asc)
            return result.head(limit) if limit else result

        if stype == "kpi_cards":
            metrics = {"Total Records": len(df)}
            if group_by and group_by in df.columns:
                metrics[f"Unique {group_by}"] = int(df[group_by].nunique())
            for col in display_cols:
                if col not in df.columns:
                    continue
                if pd.api.types.is_numeric_dtype(df[col]):
                    fn = {"sum": df[col].sum, "mean": df[col].mean,
                          "min": df[col].min, "max": df[col].max,
                          "count": df[col].count,
                          "nunique": df[col].nunique}.get(agg_func, df[col].count)
                    val = fn()
                    label = "Distinct" if agg_func == "nunique" else agg_func.title()
                    metrics[f"{label} {col}"] = round(float(val), 2) if isinstance(val, float) else int(val)
                else:
                    metrics[f"Unique {col}"] = int(df[col].nunique())
            return metrics

    except Exception as e:
        return pd.DataFrame({"Error": [str(e)]})

    return pd.DataFrame()


def build_report_from_config(df, col_config, section_configs, filter_rules, week_mode="last"):
    """Build a fully config-driven report dict."""
    filtered_df = apply_filter_rules(df, filter_rules)
    work = filtered_df.copy()

    sections = []
    for cfg in section_configs:
        if not cfg.get("enabled", True):
            continue
        data = build_section_data(work, cfg)
        sections.append({
            "id": cfg.get("id", ""),
            "title": cfg.get("title", "Section"),
            "type": cfg.get("section_type", "table"),
            "layout": cfg.get("layout", "vertical"),
            "data": data,
            "cfg": cfg,
        })

    # Build legacy keys so export functions keep working
    summary_data = next(
        (s["data"] for s in sections if s["type"] == "pivot" and isinstance(s["data"], pd.DataFrame)),
        pd.DataFrame(),
    )
    bullets = [
        f"Total records in report: {len(work)}",
        f"Sections configured: {len(sections)}",
        f"Active filter rules: {len(filter_rules)}",
    ]

    return {
        "week_range": week_range(week_mode),
        "week_mode": week_mode,
        "working": work,
        "sections": sections,
        "col_config": col_config,
        "filter_criteria": {r["col"]: f"{r['op']} {r['val']}" for r in filter_rules},
        # legacy export compatibility
        "summary": summary_data,
        "aging": pd.DataFrame(),
        "delivery": pd.DataFrame(),
        "risks": pd.DataFrame(),
        "outlook": pd.DataFrame(),
        "bullets": bullets,
    }


# ============================================================
# FILTER ENGINE
# ============================================================

def apply_prompt(df, prompt):
    work = df.copy()
    q = str(prompt).lower()
    status_col = guess_status_column(work)
    filter_dict = {}
    if status_col:
        sc = status_col
        all_statuses = work[sc].astype(str).str.lower().unique()
        all_statuses = [s for s in all_statuses if s.lower() != 'nan']
        words = q.split()
        potential_statuses = [w for w in words if len(w) > 2]
        matched_statuses = []
        for potential in potential_statuses:
            for actual_status in all_statuses:
                if potential in actual_status or actual_status in potential:
                    original = work[sc][work[sc].astype(str).str.lower() == actual_status]
                    if len(original) > 0:
                        orig_val = original.iloc[0]
                        if orig_val not in matched_statuses:
                            matched_statuses.append(orig_val)
        if matched_statuses:
            filter_dict["status"] = matched_statuses
        else:
            if "uat" in q:
                uat_statuses = [s for s in all_statuses if "uat" in s]
                if uat_statuses:
                    filter_dict["status"] = [work[sc][work[sc].astype(str).str.lower() == s].iloc[0] for s in uat_statuses]
            elif "development" in q or "dev" in q:
                dev_statuses = [s for s in all_statuses if "development" in s or "dev" in s]
                if dev_statuses:
                    filter_dict["status"] = [work[sc][work[sc].astype(str).str.lower() == s].iloc[0] for s in dev_statuses]
            elif "go live" in q or "golive" in q or "production" in q:
                prod_statuses = [s for s in all_statuses if "production" in s or "golive" in s or "go live" in s]
                if prod_statuses:
                    filter_dict["status"] = [work[sc][work[sc].astype(str).str.lower() == s].iloc[0] for s in prod_statuses]
            elif "ud" in q:
                ud_statuses = [s for s in all_statuses if "ud" in s]
                if ud_statuses:
                    filter_dict["status"] = [work[sc][work[sc].astype(str).str.lower() == s].iloc[0] for s in ud_statuses]
    nums = re.findall(r"\d+", q)
    if ("above" in q or "greater" in q or "more than" in q) and nums:
        filter_dict["aging_min"] = int(nums[0])
    if ("below" in q or "less than" in q) and nums:
        filter_dict["aging_max"] = int(nums[0])
    if "top aging" in q or "highest aging" in q:
        filter_dict["sort"] = "aging_desc"
    elif "lowest aging" in q:
        filter_dict["sort"] = "aging_asc"
    if "top" in q and nums:
        n = int(nums[0])
        if n > 0:
            filter_dict["limit"] = n
    skip_words = {
        "show", "only", "project", "projects", "above", "below", "aging", "top",
        "sort", "uat", "development", "go", "live", "golive", "greater", "than",
        "more", "less", "any", "filter", "by", "for", "with", "in"
    }
    words = [w for w in q.split() if len(w) > 3 and w not in skip_words]
    if words:
        filter_dict["search"] = " ".join(words)
    return apply_filters(work, filter_dict)


def apply_filters(df, filter_dict):
    work = df.copy()
    if work is None or len(work) == 0:
        return work
    if isinstance(filter_dict, str):
        return apply_prompt(work, filter_dict)
    if isinstance(filter_dict, dict) and "_raw" in filter_dict:
        return apply_prompt(work, filter_dict["_raw"])
    if not isinstance(filter_dict, dict):
        return work
    if "status" in filter_dict and isinstance(filter_dict["status"], list):
        sc = guess_status_column(work)
        if sc:
            status_mask = pd.Series(False, index=work.index)
            for sv in filter_dict["status"]:
                status_mask |= work[sc].astype(str).str.contains(sv, case=False, na=False)
            work = work[status_mask]
    if "project" in filter_dict and isinstance(filter_dict["project"], list):
        pc = guess_project_column(work)
        if pc:
            project_mask = pd.Series(False, index=work.index)
            for pv in filter_dict["project"]:
                project_mask |= work[pc].astype(str).str.contains(pv, case=False, na=False)
            work = work[project_mask]
        else:
            project_mask = pd.Series(False, index=work.index)
            for pv in filter_dict["project"]:
                for col in work.columns:
                    if work[col].dtype == 'object':
                        project_mask |= work[col].astype(str).str.contains(pv, case=False, na=False)
            work = work[project_mask]
    if "Aging" in work.columns:
        if "aging_min" in filter_dict:
            work = work[work["Aging"] >= filter_dict["aging_min"]]
        if "aging_max" in filter_dict:
            work = work[work["Aging"] <= filter_dict["aging_max"]]
    if "search" in filter_dict and isinstance(filter_dict["search"], str):
        term = filter_dict["search"].lower()
        if term:
            mask = work.apply(lambda row: term in " ".join([str(v).lower() for v in row.values]), axis=1)
            work = work[mask]
    if "sort" in filter_dict and "Aging" in work.columns:
        sort_type = filter_dict["sort"].lower()
        if sort_type == "aging_desc":
            work = work.sort_values("Aging", ascending=False)
        elif sort_type == "aging_asc":
            work = work.sort_values("Aging", ascending=True)
    if "limit" in filter_dict:
        limit = int(filter_dict["limit"])
        if limit > 0:
            work = work.head(limit)
    return work.reset_index(drop=True)

# ============================================================
# LLM AGENT LOGIC
# ============================================================

def classify_prompt(prompt: str) -> str:
    text = prompt.lower()
    if any(word in text for word in ['revise', 'rewrite', 'polish', 'refresh', 'reword']):
        return 'revise'
    if any(word in text for word in ['summary', 'bullet', 'executive', 'highlight', 'insight']):
        return 'summary'
    if any(word in text for word in ['status report', 'weekly report', 'wsr', 'report']):
        return 'report'
    return 'rag'


def llm_understand_prompt(prompt):
    try:
        txt = str(prompt).strip()
        resp = llm.invoke(f"""
You are a smart data analyst assistant. Convert the user's natural language request into JSON filter instructions.
Return ONLY valid JSON, no other text.

Filter fields:
- "status": array of status values (keep original case)
- "aging_min": minimum aging days (number)
- "aging_max": maximum aging days (number)
- "project": array of project names or partial names
- "sort": "aging_desc" or "aging_asc"
- "limit": number of rows (number)
- "search": text to search across all columns
- "action": "filter" or "show"

Examples:
Request: "show only UAT projects"
Response: {{"status": ["UAT"], "action": "filter"}}

Request: "top 10 aging projects"
Response: {{"sort": "aging_desc", "limit": 10, "action": "filter"}}

Request: "development projects aging more than 30 days"
Response: {{"status": ["Development"], "aging_min": 30, "action": "filter"}}

Request: "show only project ABC"
Response: {{"project": ["ABC"], "action": "filter"}}

User Request: {txt}

Return only the JSON filter instruction:
""")
        out = str(resp.content).strip()
        if "```" in out:
            out = out.split("```")[1]
            if "json" in out:
                out = out.split("json", 1)[1]
            out = out.strip()
        try:
            return json.loads(out)
        except:
            txt_lower = txt.lower()
            project_indicators = ["project", "module", "show only", "filter by", "only show"]
            if any(ind in txt_lower for ind in project_indicators):
                project_names = []
                for keyword in ["project", "module"]:
                    idx = txt_lower.find(keyword)
                    if idx != -1:
                        after = txt[idx + len(keyword):].strip(" :.,")
                        match = re.search(r'\b([A-Za-z0-9][A-Za-z0-9\-_\.]*[A-Za-z0-9])\b', after)
                        if match:
                            project_names.append(match.group(1))
                if not project_names:
                    all_matches = re.findall(r'\b([A-Za-z0-9][A-Za-z0-9\-_\.]{1,}[A-Za-z0-9])\b', txt)
                    exclude = {'only', 'show', 'filter', 'by', 'and', 'or', 'projects', 'status', 'aging', 'development', 'uat', 'testing', 'production', 'go', 'live', 'golive', 'cancelled', 'completed', 'pending'}
                    project_names = [m for m in all_matches if m.lower() not in exclude and len(m) >= 2][:3]
                if project_names:
                    return {"project": project_names, "action": "filter"}
            return {"_raw": txt, "action": "filter"}
    except Exception:
        return {"_raw": str(prompt), "action": "filter"}


def generate_from_prompt(prompt):
    src = None
    if "source_df" in st.session_state and st.session_state.source_df is not None:
        src = st.session_state.source_df.copy()
    if src is None or len(src) == 0:
        add_chat("assistant", "Please upload file and generate initial report first.")
        return None
    if "Aging" not in src.columns:
        src = prepare_data(src)
    filter_dict = llm_understand_prompt(prompt)
    if isinstance(filter_dict, dict) and "_raw" in filter_dict:
        prompt_lower = prompt.lower()
        if any(word in prompt_lower for word in ["project", "module", "show only", "filter by"]):
            project_names = []
            for keyword in ["project", "module"]:
                idx = prompt_lower.find(keyword)
                if idx != -1:
                    after = prompt[idx + len(keyword):].strip(" :.,")
                    match = re.search(r'\b([A-Za-z0-9][A-Za-z0-9\-_\.]*[A-Za-z0-9])\b', after)
                    if match and len(match.group(1)) >= 2:
                        project_names.append(match.group(1))
            if not project_names:
                all_matches = re.findall(r'\b([A-Za-z0-9][A-Za-z0-9\-_\.]{1,}[A-Za-z0-9])\b', prompt)
                exclude = {'only', 'show', 'filter', 'by', 'and', 'or', 'projects', 'status', 'aging'}
                project_names = [m for m in all_matches if m.lower() not in exclude and len(m) >= 2][:3]
            if project_names:
                filter_dict = {"project": project_names, "action": "filter"}
    filtered = apply_filters(src.copy(), filter_dict)
    if filtered is None or len(filtered) == 0:
        add_chat("assistant", "No matching records found. Showing full report.")
        filtered = src.copy()
    week_mode = st.session_state.get("week_mode", "last")
    rpt = build_report(filtered, week_mode=week_mode)
    rpt["working"] = filtered.copy()
    rpt["filter_criteria"] = filter_dict
    st.session_state.source_df = src.copy()
    st.session_state.working_df = filtered.copy()
    st.session_state.report = rpt
    st.session_state.last_prompt = prompt
    st.session_state.awaiting_next_action = True
    return rpt

# ============================================================
# PROJECT SEARCH
# ============================================================

def extract_query_from_prompt(prompt):
    p = str(prompt).strip()
    lower = p.lower()
    priority_phrases = [
        ("for", 3), ("of", 2), ("about", 5), ("on", 2), ("tell me", 8),
        ("show me", 7), ("details", 7), ("info", 4), ("status", 6), ("project", 7), ("module", 6)
    ]
    for phrase, min_len in priority_phrases:
        idx = lower.find(phrase)
        if idx != -1:
            potential = p[idx + len(phrase):].strip(" :.-,")
            tokens = re.findall(r'\b[A-Za-z0-9][A-Za-z0-9\-_\.]{2,}[A-Za-z0-9]\b', potential)
            if tokens:
                exclude = {'uat', 'dev', 'development', 'testing', 'qa', 'production', 'golive', 'go', 'live'}
                valid = [t for t in tokens if t.lower() not in exclude]
                if valid:
                    return max(valid, key=len)
    candidates = re.findall(r'\b([A-Za-z0-9][A-Za-z0-9\-_\.]*[A-Za-z0-9])\b', p)
    if candidates:
        exclude = {'uat', 'dev', 'development', 'testing', 'qa', 'production', 'golive', 'go', 'live', 'status', 'project', 'module', 'details'}
        candidates = [c for c in candidates if len(c) >= 3 and c.lower() not in exclude]
        if candidates:
            return max(candidates, key=len)
    return p


def find_project_rows(df, query):
    q = str(query).strip().lower()
    if not q:
        return pd.DataFrame()
    q = re.sub(r'\s+', ' ', q).strip()
    priority_cols = ["project", "module", "project name", "module name", "id", "module id", "project id", "code"]
    candidate_cols = [col for col in df.columns if any(pc in col.lower() for pc in priority_cols)]
    if not candidate_cols:
        candidate_cols = list(df.columns)
    mask = pd.Series(False, index=df.index)
    for col in candidate_cols[:5]:
        try:
            col_values = df[col].astype(str).str.lower()
            exact_match = col_values == q
            mask |= exact_match
            if exact_match.any():
                break
            partial_match = col_values.str.contains(re.escape(q), case=False, regex=True, na=False)
            mask |= partial_match
        except:
            pass
    if mask.any():
        return df[mask].drop_duplicates()
    mask = pd.Series(False, index=df.index)
    for col in candidate_cols:
        try:
            col_values = df[col].astype(str).str.lower()
            word_match = col_values.str.contains(r'\b' + re.escape(q) + r'\b', case=False, regex=True, na=False)
            mask |= word_match
        except:
            pass
    if mask.any():
        return df[mask].drop_duplicates()
    try:
        full_mask = df.apply(lambda row: q in " ".join([str(v).lower() for v in row.values]), axis=1)
        if full_mask.any():
            return df[full_mask].drop_duplicates().head(10)
    except:
        pass
    return pd.DataFrame()


def format_project_details(rows, df):
    if rows is None or len(rows) == 0:
        return "No matching project found. Please try a more complete project name or module ID."
    output = []
    cols = std_cols(df)
    project_col = cols.get("project") or find_col(df, ["Project", "Project Name", "Project Module", "Module"])
    status_col = cols.get("status") or find_col(df, ["Status", "Status Code", "Current Status"])
    uat_col = find_col(df, ["Actual UAT Start", "UAT Start", "UAT Date"])
    golive_col = find_col(df, ["Actual Go live", "Actual Go Live", "Go Live Date"])
    desc_col = find_col(df, ["Description", "Comments", "Latest Comment", "Remarks"])
    date_col = find_col(df, ["Latest Date", "Updated Date", "Modified Date", "Comment Date", "Last Updated"])
    module_id_col = find_col(df, ["Module ID", "Module Code", "Project ID", "Project Code", "ID"])
    aging_col = "Aging" if "Aging" in df.columns else None
    for i, (_, row) in enumerate(rows.head(5).iterrows()):
        project_name = safe_str(row.get(project_col, "")) if project_col else "Unknown"
        module_id = safe_str(row.get(module_id_col, "")) if module_id_col else ""
        status = safe_str(row.get(status_col, "")) if status_col else ""
        aging = safe_str(row.get(aging_col, "")) if aging_col else ""
        aging_from = safe_str(row.get("Aging From Date", ""))
        uat_date = safe_str(row.get(uat_col, "")) if uat_col else ""
        golive_date = safe_str(row.get(golive_col, "")) if golive_col else ""
        description = safe_str(row.get(desc_col, "")) if desc_col else ""
        latest_update = safe_str(row.get(date_col, "")) if date_col else ""
        if i == 0:
            output.append("## Project Details")
        output.append(f"\n### {project_name}")
        if module_id:
            output.append(f"**Module ID:** {module_id}")
        status_info = []
        if status:
            emoji = "🟢"
            if "hold" in status.lower() or "delay" in status.lower():
                emoji = "🟡"
            elif "cancel" in status.lower() or "stop" in status.lower():
                emoji = "🔴"
            elif "complete" in status.lower():
                emoji = "✅"
            status_info.append(f"{emoji} **Status:** {status}")
        if aging:
            aging_days = int(float(aging)) if aging.replace('.', '').isdigit() else 0
            emoji = "🟢" if aging_days < 30 else "🟡" if aging_days < 60 else "🔴"
            status_info.append(f"{emoji} **Aging:** {aging} days")
        if status_info:
            output.append("\n**Current Status:**")
            output.extend([f"- {info}" for info in status_info])
        timeline_info = []
        if aging_from:
            timeline_info.append(f"📅 **Aging Start:** {aging_from}")
        if uat_date:
            timeline_info.append(f"🧪 **UAT Date:** {uat_date}")
        if golive_date:
            timeline_info.append(f"🚀 **Go Live Date:** {golive_date}")
        if latest_update:
            timeline_info.append(f"📝 **Last Updated:** {latest_update}")
        if timeline_info:
            output.append("\n**Timeline:**")
            output.extend([f"- {info}" for info in timeline_info])
        if description:
            output.append(f"\n**Description:** {description}")
        if len(rows) > 1 and i < min(5, len(rows)) - 1:
            output.append("\n---")
    if len(rows) > 5:
        output.append(f"\n⚠️ {len(rows)} matches found. Showing first 5.")
    return "\n".join(output)


def is_project_details_request(prompt):
    p = str(prompt).strip().lower()
    project_intent = ["detail", "details", "tell me", "info", "information", "what is", "what are", "status of", "about", "find", "get info", "project info", "module info", "show details"]
    filter_keywords = ["show only", "filter", "only show", "display only", "list only", "filter by", "give me", "list", "filter to", "display", "view only"]
    status_indicators = ["uat", "development", "dev", "go live", "golive", "production", "in progress", "on hold", "cancelled", "completed", "pending", "ud submitted", "ud", "testing", "qa", "review", "signedoff", "uat signedoff", "sit", "preuat"]
    if any(k in p for k in filter_keywords):
        return False
    has_status = any(s in p for s in status_indicators)
    has_filter_intent = any(k in p for k in ["only", "filter", "show", "display", "list"])
    if has_status and has_filter_intent:
        return False
    has_project_intent = any(k in p for k in project_intent)
    potential_projects = re.findall(r'\b[A-Za-z0-9][A-Za-z0-9\-_\.]{2,}[A-Za-z0-9]\b', p)
    actual_projects = [proj for proj in potential_projects if proj.lower() not in status_indicators]
    has_specific_project = len(actual_projects) > 0
    has_project_keywords = any(k in p for k in ["project", "module", "id ", " id", "code", "module id"])
    return (has_project_intent and has_specific_project) or (has_project_keywords and has_specific_project)


def project_details_response(prompt):
    src = None
    if st.session_state.source_df is not None:
        src = st.session_state.source_df.copy()
    if src is None or len(src) == 0:
        return "Please upload data and generate the initial report first."
    if "Aging" not in src.columns:
        src = prepare_data(src)
    query = extract_query_from_prompt(prompt)
    rows = find_project_rows(src, query)
    return format_project_details(rows, src)

# ============================================================
# INTENT DETECTION
# ============================================================

def wants_email(prompt):
    p = str(prompt).lower()
    return any(k in p for k in ["email", "mail", "send report", "send mail", "share report"])


def wants_modify(prompt):
    p = str(prompt).lower()
    return any(k in p for k in ["modify", "change", "filter", "remove", "add", "rename", "update", "regenerate", "show only"])


def wants_export(prompt):
    p = str(prompt).lower()
    return any(k in p for k in ["export", "download", "save", "get pdf", "get ppt", "get image", "get excel", "generate pdf", "generate ppt", "generate image", "generate excel"])


def get_export_format(prompt):
    p = str(prompt).lower()
    if "pdf" in p:
        return "pdf"
    elif "ppt" in p or "powerpoint" in p:
        return "ppt"
    elif "image" in p or "png" in p:
        return "image"
    elif "excel" in p or "xlsx" in p:
        return "excel"
    return None


def parse_emails(txt):
    return re.findall(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", txt)


def parse_week_mode(prompt):
    p = str(prompt).lower()
    if "current week" in p or "this week" in p:
        return "current"
    if "last week" in p or "previous week" in p:
        return "last"
    return None


def is_week_report_request(prompt):
    p = str(prompt).lower()
    mode = parse_week_mode(p)
    if mode is None:
        return False
    return any(k in p for k in ["report", "generate", "build", "show", "create"])


def assistant_reply(prompt, rpt):
    p = norm(prompt)
    if rpt is None:
        return "Please upload a source file first."
    if "summary" in p:
        return "Report summary regenerated successfully."
    if "uat" in p:
        return "Filtered UAT portfolio report generated."
    if "go live" in p:
        return "Go Live focused report generated."
    if "risk" in p:
        return "Risk focused report generated."
    if "project" in p:
        return "Project specific report generated."
    return "Report updated successfully."


def ask_gemini_direct(prompt):
    try:
        response = llm.invoke([HumanMessage(content=prompt)])
        return response.content
    except:
        return "Unable to process request."

# ============================================================
# VALIDATION
# ============================================================

def validate_uploaded_data(df):
    results = {"errors": [], "warnings": [], "info": []}
    required_patterns = {
        "project": ["project", "module", "initiative"],
        "status": ["status", "state", "phase"],
        "date": ["date", "uat", "go live", "start", "end"]
    }
    found_columns = {}
    for category, patterns in required_patterns.items():
        found = False
        for pattern in patterns:
            for col in df.columns:
                if pattern.lower() in col.lower():
                    found_columns[category] = col
                    found = True
                    break
            if found:
                break
        if not found:
            results["errors"].append(f"Missing {category} column. Expected patterns: {patterns}")
    if len(df) == 0:
        results["errors"].append("File appears to be empty")
    elif len(df) < 5:
        results["warnings"].append("Very small dataset detected.")
    if "project" in found_columns:
        duplicates = df[found_columns["project"]].duplicated().sum()
        if duplicates > 0:
            results["warnings"].append(f"Found {duplicates} duplicate project entries")
    if results["errors"]:
        results["info"].append(f"❌ {len(results['errors'])} critical issues found")
    if results["warnings"]:
        results["info"].append(f"⚠️ {len(results['warnings'])} warnings")
    if not results["errors"] and not results["warnings"]:
        results["info"].append("✅ Data validation passed")
    return results


def display_validation_results(results):
    if results["errors"]:
        for error in results["errors"]:
            st.error(f"• {error}")
    if results["warnings"]:
        for warning in results["warnings"]:
            st.warning(f"• {warning}")
    if results["info"]:
        for info in results["info"]:
            if "✅" in info:
                st.success(info)
            elif "❌" in info:
                st.error(info)
            elif "⚠️" in info:
                st.warning(info)
            else:
                st.info(info)

# ============================================================
# PDF EXPORT HELPERS
# ============================================================

def pdf_section(title):
    styles = getSampleStyleSheet()
    section_style = ParagraphStyle(
        "section", parent=styles["Heading3"],
        fontName="Helvetica-Bold", fontSize=11,
        textColor=colors.HexColor("#ff00ff"),
        alignment=0, spaceBefore=8, spaceAfter=4,
    )
    return Paragraph(f"<b>{title}</b>", section_style)


def pdf_table(df, widths):
    if df is None or len(df) == 0:
        return Paragraph("No data available", getSampleStyleSheet()["Normal"])
    temp = df.copy()
    styles = getSampleStyleSheet()
    normal_style = ParagraphStyle(
        "cell", parent=styles["BodyText"],
        fontName="Helvetica", fontSize=8, leading=10, alignment=0,
        leftIndent=0, rightIndent=0, spaceBefore=0, spaceAfter=0,
    )
    header_style = ParagraphStyle(
        "header", parent=styles["BodyText"],
        fontName="Helvetica-Bold", fontSize=8, leading=10,
        textColor=colors.white, alignment=0,
        leftIndent=0, rightIndent=0, spaceBefore=0, spaceAfter=0,
    )
    if len(temp.columns) > 1:
        wide_cols = [
            3.2 * inch if "project" in str(c).lower() or "module" in str(c).lower() else None
            for c in temp.columns
        ]
        if any(w is not None for w in wide_cols):
            remaining = 10.8 * inch - sum(w for w in wide_cols if w is not None)
            other_count = sum(1 for w in wide_cols if w is None)
            if other_count > 0 and remaining > 0:
                widths = [w if w is not None else remaining / other_count for w in wide_cols]
    rows = [[Paragraph(str(c), header_style) for c in temp.columns]]
    for _, r in temp.iterrows():
        rows.append([Paragraph(str(v), normal_style) for v in r.tolist()])
    t = Table(rows, colWidths=widths, repeatRows=1)
    style = [
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#01469a")),
        ("GRID", (0, 0), (-1, -1), 0.6, CLR_BORDER),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ("WORDWRAP", (0, 0), (-1, -1), True),
    ]
    for i in range(1, len(rows)):
        txt = " ".join([str(cell).lower() for cell in rows[i]])
        if "grand total" in txt:
            style += [
                ("BACKGROUND", (0, i), (-1, i), CLR_TOTAL),
                ("FONTNAME", (0, i), (-1, i), "Helvetica-Bold")
            ]
    t.setStyle(TableStyle(style))
    return t


def format_ppt_table(tbl, header_color=RGBColor(1, 70, 150), header_font_size=10):
    from pptx.util import Pt
    from pptx.enum.text import MSO_ANCHOR
    for cell in tbl.rows[0].cells:
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = header_color
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(header_font_size)
                run.font.color.rgb = RGBColor(255, 255, 255)
    for row in tbl.rows:
        for cell in row.cells:
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size is None:
                        run.font.size = Pt(9)

# ============================================================
# CHART HELPERS
# ============================================================

def export_status_chart_image(report, folder):
    try:
        working_df = report["working"].copy()
        cols = std_cols(working_df)
        status_col = cols.get("status")
        if status_col and status_col in working_df.columns:
            status_counts = working_df[status_col].astype(str).str.lower().value_counts().reset_index()
            status_counts.columns = ["Status", "Count"]
            all_statuses = sorted(working_df[status_col].astype(str).str.lower().unique())
            all_status_df = pd.DataFrame({"Status": all_statuses})
            status_counts = all_status_df.merge(status_counts, on="Status", how="left").fillna(0)
            status_counts["Count"] = status_counts["Count"].astype(int)
            status_counts = status_counts.sort_values(by=["Count", "Status"], ascending=[False, True])
            if len(status_counts) > 0:
                fig_bar = px.bar(
                    status_counts, x="Status", y="Count",
                    title="Project Status Distribution", color="Status",
                    color_discrete_sequence=px.colors.qualitative.Plotly,
                    labels={"Count": "Number of Projects", "Status": "Project Status"}
                )
                fig_bar.update_layout(height=350, showlegend=True, xaxis_tickangle=-45, margin=dict(b=120))
                img_path = folder / "status_chart.png"
                fig_bar.write_image(str(img_path), width=800, height=350)
                return img_path
    except Exception as e:
        print(f"Error creating status chart: {e}")
    return None


def export_aging_chart_image(report, folder):
    try:
        aging_df = report["aging"].reset_index()
        aging_cols = ["0 - 30 Days", "30 - 60 Days", "60 - 90 Days", "> 90 Days"]
        aging_data = [{"Bucket": col, "Count": aging_df[col].sum()} for col in aging_cols if col in aging_df.columns]
        aging_data = [d for d in aging_data if d["Count"] > 0]
        if aging_data:
            aging_chart_df = pd.DataFrame(aging_data)
            fig_pie = px.pie(
                aging_chart_df, names="Bucket", values="Count",
                title="Project Aging Distribution",
                color_discrete_sequence=px.colors.qualitative.Set3
            )
            fig_pie.update_layout(height=350)
            img_path = folder / "aging_chart.png"
            fig_pie.write_image(str(img_path), width=500, height=350)
            return img_path
    except Exception as e:
        print(f"Error creating aging chart: {e}")
    return None

# ============================================================
# EXPORT: PDF
# ============================================================

def export_pdf(report):
    folder = output_folder()
    file = folder / "WSR_Report.pdf"
    doc = SimpleDocTemplate(
        str(file), pagesize=landscape(A4),
        leftMargin=14, rightMargin=14, topMargin=14, bottomMargin=14
    )
    styles = getSampleStyleSheet()
    story = []
    title_style = ParagraphStyle(
        "ttl", parent=styles["Title"], alignment=1, fontSize=18, fontName="Helvetica-Bold"
    )
    story.append(Paragraph(f"<u>{get_report_title(report)}</u>", title_style))
    story.append(Spacer(1, 10))
    story.append(pdf_section("1. Project Status"))
    story.append(Spacer(1, 5))
    story.append(pdf_table(report["summary"], [10.8 * inch / max(1, len(report["summary"].columns))] * len(report["summary"].columns)))
    story.append(Spacer(1, 8))
    story.append(pdf_section("2. Project Aging Summary"))
    story.append(Spacer(1, 5))
    ag = report["aging"].reset_index()
    story.append(pdf_table(ag, [10.8 * inch / max(1, len(ag.columns))] * len(ag.columns)))
    story.append(Spacer(1, 8))
    story.append(pdf_section("3. Delivery Highlights"))
    story.append(Spacer(1, 5))
    dl = report["delivery"]
    if len(dl) == 0:
        dl = pd.DataFrame({"Message": ["No delivery in last week"]})
    story.append(pdf_table(dl, [10.8 * inch / max(1, len(dl.columns))] * len(dl.columns)))
    story.append(Spacer(1, 8))
    story.append(pdf_section("4. Risks / Aging Projects"))
    story.append(Spacer(1, 5))
    rk = report["risks"]
    story.append(pdf_table(rk, [10.8 * inch / max(1, len(rk.columns))] * len(rk.columns)))
    story.append(Spacer(1, 8))
    story.append(pdf_section("5. Executive Summary"))
    story.append(Spacer(1, 6))
    bullet_style = ParagraphStyle(
        "bullet", parent=styles["Normal"],
        fontSize=10, leading=14, leftIndent=10, spaceBefore=4, spaceAfter=4
    )
    for b in report["bullets"]:
        story.append(Paragraph(f"• {b}", bullet_style))
        story.append(Spacer(1, 4))
    story.append(Spacer(1, 8))
    story.append(pdf_section("6. Outlook for Next Week"))
    story.append(Spacer(1, 5))
    ol = report["outlook"]
    if len(ol) == 0:
        ol = pd.DataFrame({"Message": ["No deliveries scheduled for next week"]})
    story.append(pdf_table(ol, [10.8 * inch / max(1, len(ol.columns))] * len(ol.columns)))
    story.append(Spacer(1, 8))
    story.append(pdf_section("7. Status Dashboard"))
    story.append(Spacer(1, 5))
    status_chart_path = export_status_chart_image(report, folder)
    if status_chart_path and status_chart_path.exists():
        from reportlab.platypus import Image as RLImage
        story.append(RLImage(str(status_chart_path), width=6 * inch, height=2.5 * inch))
        story.append(Spacer(1, 8))
    aging_chart_path = export_aging_chart_image(report, folder)
    if aging_chart_path and aging_chart_path.exists():
        from reportlab.platypus import Image as RLImage
        story.append(RLImage(str(aging_chart_path), width=4 * inch, height=2.5 * inch))
        story.append(Spacer(1, 8))
    doc.build(story)
    return file

# ============================================================
# EXPORT: PPT
# ============================================================

def export_ppt(report):
    folder = output_folder()
    file = folder / "WSR_Report.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(10)

    def add_slide():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def add_title(slide, title_text):
        tx = slide.shapes.add_textbox(Inches(.2), Inches(.1), Inches(12.8), Inches(.4))
        p = tx.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title_text
        r.font.bold = True
        r.font.underline = True
        r.font.size = Pt(22)
        p.alignment = PP_ALIGN.CENTER

    def add_section_label(slide, label, y):
        tb = slide.shapes.add_textbox(Inches(.2), Inches(y), Inches(12.9), Inches(.25))
        tb.text_frame.text = label
        for para in tb.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = RGBColor(255, 0, 255)
                run.font.bold = True

    def add_df_table(slide, df, y, max_h=2.5):
        rows, cols = len(df) + 1, len(df.columns)
        tbl = slide.shapes.add_table(rows, cols, Inches(.2), Inches(y), Inches(12.9), Inches(min(max_h, 9 - y))).table
        for c_idx, col in enumerate(df.columns):
            tbl.cell(0, c_idx).text = str(col)
        for r_idx in range(1, rows):
            for c_idx in range(cols):
                tbl.cell(r_idx, c_idx).text = str(df.iloc[r_idx - 1, c_idx])
        format_ppt_table(tbl)

    # Slide 1: Title + Project Status + Aging
    s1 = add_slide()
    add_title(s1, get_report_title(report))
    y = 0.65
    add_section_label(s1, "1. Project Status", y)
    y += 0.3
    sm = report["summary"]
    add_df_table(s1, sm, y, max_h=0.8)
    y += 1.1
    add_section_label(s1, "2. Project Aging Summary", y)
    y += 0.3
    ag = report["aging"].reset_index()
    add_df_table(s1, ag, y, max_h=2.5)

    # Slide 2: Delivery + Risks
    s2 = add_slide()
    add_title(s2, get_report_title(report))
    y = 0.65
    add_section_label(s2, "3. Delivery Highlights", y)
    y += 0.3
    dl = report["delivery"] if len(report["delivery"]) > 0 else pd.DataFrame({"Message": ["No delivery in last week"]})
    add_df_table(s2, dl, y, max_h=2.5)
    y += min(2.8, 0.3 + len(dl) * 0.35)
    add_section_label(s2, "4. Risks / Aging Projects", y)
    y += 0.3
    rk = report["risks"] if len(report["risks"]) > 0 else pd.DataFrame({"Message": ["No risks identified"]})
    add_df_table(s2, rk, y, max_h=2.5)

    # Slide 3: Executive Summary + Outlook
    s3 = add_slide()
    add_title(s3, get_report_title(report))
    y = 0.65
    add_section_label(s3, "5. Executive Summary", y)
    y += 0.3
    for b in report["bullets"]:
        tx = s3.shapes.add_textbox(Inches(.12), Inches(y), Inches(12.8), Inches(.75))
        tx.text_frame.word_wrap = True
        tx.text_frame.text = "• " + str(b)
        for para in tx.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(12)
        y += 0.95
    y += 0.3
    add_section_label(s3, "6. Outlook for Next Week", y)
    y += 0.3
    ol = report["outlook"] if len(report["outlook"]) > 0 else pd.DataFrame({"Message": ["No deliveries scheduled for next week"]})
    add_df_table(s3, ol, y, max_h=2.0)

    prs.save(str(file))
    return file

# ============================================================
# EXPORT: IMAGE (PNG)
# ============================================================

def export_image(report):
    folder = output_folder()
    file = folder / "WSR_Report.png"
    ag = report["aging"].reset_index()
    dl = report["delivery"]
    if len(dl) == 0:
        dl = pd.DataFrame({"Message": ["No delivery in last week"]})
    rk = report["risks"]
    ol = report["outlook"]
    if len(ol) == 0:
        ol = pd.DataFrame({"Message": ["No deliveries scheduled for next week"]})
    height = (
        400 + 80 + 60 + 40
        + 80 + max(100, (len(ag) + 1) * 25) + 40
        + 80 + max(100, (len(dl) + 1) * 25) + 40
        + 80 + max(100, (len(rk) + 1) * 25) + 40
        + 80 + max(100, (len(ol) + 1) * 25) + 40
        + 80 + (len(report["bullets"]) * 70) + 40
    )
    img = Image.new("RGB", (1900, height), "white")
    d = ImageDraw.Draw(img)
    try:
        title_font = ImageFont.truetype("arial.ttf", 20)
        section_font = ImageFont.truetype("arial.ttf", 16)
        header_font = ImageFont.truetype("arial.ttf", 14)
        text_font = ImageFont.truetype("arial.ttf", 12)
    except:
        title_font = section_font = header_font = text_font = ImageFont.load_default()

    def draw_table(data_df, y_start, cols_list, col_width, aging_mode=False):
        y = y_start
        for col_idx, col in enumerate(cols_list):
            x = 20 + (col_idx * col_width)
            col_name = str(col).strip().lower()
            fill_color = (1, 70, 150)
            if aging_mode:
                if col_name == "0 - 30 days":
                    fill_color = (22, 163, 74)
                elif col_name in ["30 - 60 days", "60 - 90 days"]:
                    fill_color = (245, 158, 11)
                elif col_name in [">90 days", "> 90 days"]:
                    fill_color = (220, 38, 38)
            d.rectangle((x, y, x + col_width, y + 25), fill=fill_color)
            d.rectangle((x, y, x + col_width, y + 25), outline=(100, 100, 100), width=1)
            d.text((x + 5, y + 5), str(col), fill="white", font=header_font)
        y += 25
        for row_idx, row in data_df.iterrows():
            row_text = " ".join([str(v).lower() for v in row.tolist()])
            if "grand total" in row_text:
                row_bg = (255, 243, 205)
            else:
                row_bg = (243, 243, 243) if row_idx % 2 == 0 else (255, 255, 255)
            for col_idx, val in enumerate(row):
                x = 20 + (col_idx * col_width)
                d.rectangle((x, y, x + col_width, y + 25), fill=row_bg)
                d.rectangle((x, y, x + col_width, y + 25), outline=(180, 180, 180), width=1)
                d.text((x + 5, y + 5), str(val), fill="black", font=text_font)
            y += 25
        return y

    d.rectangle((0, 0, 1900, 60), fill=(46, 125, 50))
    title_text = get_report_title(report)
    try:
        title_bbox = d.textbbox((0, 0), title_text, font=title_font)
        title_width = title_bbox[2] - title_bbox[0]
    except:
        title_width = 600
    title_x = (1900 - title_width) // 2
    d.text((title_x, 12), title_text, fill="white", font=title_font)
    d.line((title_x, 50, title_x + title_width, 50), fill="white", width=2)
    y = 80
    d.text((30, y + 8), "1. Project Status", fill=(255, 0, 255), font=section_font)
    y += 45
    sm = report["summary"]
    cw = 1840 // max(1, len(sm.columns))
    y = draw_table(sm, y, list(sm.columns), cw)
    y += 20
    d.text((30, y + 8), "2. Project Aging Summary", fill=(255, 0, 255), font=section_font)
    y += 45
    cw = 1840 // max(1, len(ag.columns))
    y = draw_table(ag, y, list(ag.columns), cw, aging_mode=True)
    y += 20
    d.text((30, y + 8), "3. Delivery Highlights", fill=(255, 0, 255), font=section_font)
    y += 45
    cw = 1840 // max(1, len(dl.columns))
    y = draw_table(dl, y, list(dl.columns), cw)
    y += 20
    d.text((30, y + 8), "4. Risks / Aging Projects", fill=(255, 0, 255), font=section_font)
    y += 45
    cw = 1840 // max(1, len(rk.columns))
    y = draw_table(rk, y, list(rk.columns), cw) if len(rk) > 0 else y + 25
    y += 40
    d.text((30, y + 8), "5. Outlook for Next Week", fill=(255, 0, 255), font=section_font)
    y += 45
    cw = 1840 // max(1, len(ol.columns))
    y = draw_table(ol, y, list(ol.columns), cw)
    y += 40
    d.text((30, y + 8), "6. Executive Summary", fill=(255, 0, 255), font=section_font)
    y += 45
    for bullet in report["bullets"]:
        d.text((30, y + 15), "• " + str(bullet), fill="black", font=text_font)
        y += 40
    img.save(str(file))
    return file

# ============================================================
# EXPORT: EXCEL
# ============================================================

def export_excel(report):
    folder = output_folder()
    file = folder / "WSR_Final.xlsx"
    final_df = report["working"].copy()
    for c in final_df.columns:
        if is_date_col(c):
            final_df[c] = final_df[c].apply(fmt_date)
    with pd.ExcelWriter(file, engine="openpyxl") as writer:
        final_df.to_excel(writer, sheet_name="Final Data", index=False)
    return file

# ============================================================
# EMAIL
# ============================================================

def send_email(recipients, report):
    host = st.session_state.smtp_host
    port = int(st.session_state.smtp_port)
    user = st.session_state.smtp_user
    pwd = st.session_state.smtp_pass
    sender = st.session_state.smtp_sender or user
    if not host or not user or not pwd:
        return False, "❌ SMTP settings missing. Please configure Email Settings."
    if not recipients or len(recipients) == 0:
        return False, "❌ No recipient email addresses provided."
    try:
        img = export_image(report)
        xl = export_excel(report)
    except Exception as e:
        return False, f"❌ Failed to export report files: {str(e)}"
    try:
        img_path = Path(img) if not isinstance(img, Path) else img
        xl_path = Path(xl) if not isinstance(xl, Path) else xl
        if not img_path.exists():
            return False, f"❌ Report image not found: {img_path}"
        if not xl_path.exists():
            return False, f"❌ Report Excel file not found: {xl_path}"
    except Exception as e:
        return False, f"❌ File path validation failed: {str(e)}"
    try:
        msg = MIMEMultipart("related")
        msg["Subject"] = f"Weekly Status Report | {get_report_week_range(report)}"
        msg["From"] = sender
        msg["To"] = ", ".join(recipients)
        html = """
        <html><body>
            <p>Hi Team,</p>
            <p>Please find the Weekly Status Report attached.</p>
            <img src="cid:wsrimg" style="max-width: 100%; height: auto;">
            <p>Regards,<br>Smart WSR Agent</p>
        </body></html>
        """
        msg.attach(MIMEText(html, "html"))
        with open(img_path, "rb") as f:
            im = MIMEImage(f.read())
            im.add_header("Content-ID", "<wsrimg>")
            im.add_header("Content-Disposition", "inline", filename=img_path.name)
            msg.attach(im)
        with open(xl_path, "rb") as f:
            part = MIMEApplication(f.read())
            part.add_header("Content-Disposition", "attachment", filename=xl_path.name)
            msg.attach(part)
    except Exception as e:
        return False, f"❌ Failed to compose email: {str(e)}"
    server = None
    try:
        server = smtplib.SMTP(host, port, timeout=10)
        server.starttls()
        server.login(user, pwd)
        server.sendmail(sender, recipients, msg.as_string())
        server.quit()
        return True, f"✅ Email sent successfully to {len(recipients)} recipient(s)."
    except smtplib.SMTPAuthenticationError as e:
        return False, f"❌ Authentication Failed. Use App Password if 2FA enabled. Error: {str(e)[:100]}"
    except smtplib.SMTPException as e:
        if "539" in str(e):
            return False, "❌ SMTP Error 539: Authentication policy issue."
        elif "535" in str(e):
            return False, "❌ SMTP Error 535: Authentication unsuccessful."
        return False, f"❌ SMTP Error: {str(e)[:150]}"
    except (OSError, socket.error) as e:
        return False, f"❌ Connection Failed: Cannot reach {host}:{port}. Error: {str(e)[:80]}"
    except Exception as e:
        return False, f"❌ Unexpected error: {str(e)[:150]}"
    finally:
        if server:
            try:
                server.quit()
            except:
                pass


def test_smtp_connection():
    host = st.session_state.smtp_host
    port = int(st.session_state.smtp_port) if st.session_state.smtp_port else 587
    user = st.session_state.smtp_user
    pwd = st.session_state.smtp_pass
    results = {"host": host, "port": port, "user": user, "steps": []}
    if not host or not user or not pwd:
        results["status"] = "FAILED"
        results["steps"].append("❌ STEP 1 - Settings Validation: MISSING credentials")
        return results
    results["steps"].append("✅ STEP 1 - Settings Validation: PASSED")
    try:
        server = smtplib.SMTP(host, port, timeout=10)
        results["steps"].append(f"✅ STEP 2 - Network Connection: PASSED (Connected to {host}:{port})")
    except (OSError, socket.error) as e:
        results["status"] = "FAILED"
        results["steps"].append(f"❌ STEP 2 - Network Connection: FAILED - Cannot reach {host}:{port}")
        results["error"] = str(e)[:100]
        return results
    except Exception as e:
        results["status"] = "FAILED"
        results["steps"].append(f"❌ STEP 2 - Network Connection: FAILED - {str(e)[:80]}")
        return results
    try:
        server.starttls()
        results["steps"].append("✅ STEP 3 - STARTTLS Upgrade: PASSED")
    except Exception as e:
        results["status"] = "FAILED"
        results["steps"].append(f"❌ STEP 3 - STARTTLS Upgrade: FAILED - {str(e)[:80]}")
        try:
            server.quit()
        except:
            pass
        return results
    try:
        server.login(user, pwd)
        results["steps"].append("✅ STEP 4 - Authentication: PASSED")
        results["status"] = "SUCCESS"
    except smtplib.SMTPAuthenticationError as e:
        results["status"] = "FAILED"
        error_msg = str(e)
        if "539" in error_msg:
            results["steps"].append("❌ STEP 4 - Authentication: FAILED (Error 539: Policy conflict)")
        elif "535" in error_msg:
            results["steps"].append("❌ STEP 4 - Authentication: FAILED (Error 535: Credentials rejected)")
        else:
            results["steps"].append(f"❌ STEP 4 - Authentication: FAILED - {error_msg[:80]}")
        results["error"] = error_msg[:150]
    except Exception as e:
        results["status"] = "FAILED"
        results["steps"].append(f"❌ STEP 4 - Authentication: FAILED - {str(e)[:80]}")
        results["error"] = str(e)[:150]
    try:
        server.quit()
    except:
        pass
    return results
