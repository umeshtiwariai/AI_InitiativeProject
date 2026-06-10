import os
from reportlab.lib.pagesizes import A4, landscape
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
from datetime import datetime

BASE=r'D:/WSR_History'


def folder():
    w=datetime.now().isocalendar().week
    p=os.path.join(BASE,str(datetime.now().year),f'Week_{w}')
    os.makedirs(p, exist_ok=True)
    return p


def export_pdf(report):
    """Export comprehensive 5-section report to PDF"""
    f=os.path.join(folder(),'WSR_Report.pdf')
    doc=SimpleDocTemplate(f, pagesize=landscape(A4), rightMargin=30, leftMargin=30)
    s=getSampleStyleSheet()
    story=[]
    
    # Title
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=s['Heading1'],
        fontSize=18,
        textColor=colors.HexColor('#2e7d32'),
        spaceAfter=20,
        alignment=1
    )
    week_range = report.get('week_range', 'Current Week')
    story.append(Paragraph(f'Weekly Status Report : {week_range}', title_style))
    story.append(Spacer(1, 12))
    
    # Section 1: Project Status (Summary)
    story.append(Paragraph('Project Status', s['Heading2']))
    story.append(Spacer(1, 8))
    if not report['summary'].empty:
        summary_data = [list(report['summary'].columns)] + report['summary'].values.tolist()
        summary_table = Table(summary_data, colWidths=[1.2*inch]*len(report['summary'].columns))
        summary_table.setStyle(TableStyle([
            ('BACKGROUND',(0,0),(-1,0),colors.HexColor('#2e7d32')),
            ('TEXTCOLOR',(0,0),(-1,0),colors.whitesmoke),
            ('ALIGN',(0,0),(-1,-1),'CENTER'),
            ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
            ('FONTSIZE',(0,0),(-1,0),11),
            ('BOTTOMPADDING',(0,0),(-1,0),12),
            ('GRID',(0,0),(-1,-1),1,colors.grey),
            ('FONTSIZE',(0,1),(-1,-1),10),
            ('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white, colors.HexColor('#f3f3f3')])
        ]))
        story.append(summary_table)
    story.append(Spacer(1, 16))
    
    # Section 2: Project Aging Summary
    story.append(Paragraph('Project Aging Summary (Top 10)', s['Heading2']))
    story.append(Spacer(1, 8))
    if not report['aging'].empty:
        aging_data = [list(report['aging'].columns)] + report['aging'].values.tolist()
        aging_table = Table(aging_data, colWidths=[4*inch, 1*inch])
        aging_table.setStyle(TableStyle([
            ('BACKGROUND',(0,0),(-1,0),colors.HexColor('#2e7d32')),
            ('TEXTCOLOR',(0,0),(-1,0),colors.whitesmoke),
            ('ALIGN',(0,0),(-1,-1),'LEFT'),
            ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
            ('FONTSIZE',(0,0),(-1,0),10),
            ('GRID',(0,0),(-1,-1),1,colors.grey),
            ('FONTSIZE',(0,1),(-1,-1),9),
            ('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white, colors.HexColor('#f3f3f3')])
        ]))
        story.append(aging_table)
    story.append(Spacer(1, 16))
    
    # Section 3: Delivery Highlights
    story.append(Paragraph('Delivery Highlights (Last 7 Days)', s['Heading2']))
    story.append(Spacer(1, 8))
    if not report['delivery'].empty:
        delivery_data = [list(report['delivery'].columns)] + report['delivery'].values.tolist()
        delivery_table = Table(delivery_data)
        delivery_table.setStyle(TableStyle([
            ('BACKGROUND',(0,0),(-1,0),colors.HexColor('#2e7d32')),
            ('TEXTCOLOR',(0,0),(-1,0),colors.whitesmoke),
            ('ALIGN',(0,0),(-1,-1),'LEFT'),
            ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
            ('FONTSIZE',(0,0),(-1,0),10),
            ('GRID',(0,0),(-1,-1),1,colors.grey),
            ('FONTSIZE',(0,1),(-1,-1),9),
            ('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white, colors.HexColor('#f3f3f3')])
        ]))
        story.append(delivery_table)
    else:
        story.append(Paragraph('No delivery highlights in the last 7 days', s['Normal']))
    story.append(Spacer(1, 16))
    
    # Section 4: Top Risk Projects
    story.append(Paragraph('Top Risk Projects (Highest Aging)', s['Heading2']))
    story.append(Spacer(1, 8))
    if not report['risks'].empty:
        risks_data = [list(report['risks'].columns)] + report['risks'].values.tolist()
        risks_table = Table(risks_data)
        risks_table.setStyle(TableStyle([
            ('BACKGROUND',(0,0),(-1,0),colors.HexColor('#2e7d32')),
            ('TEXTCOLOR',(0,0),(-1,0),colors.whitesmoke),
            ('ALIGN',(0,0),(-1,-1),'LEFT'),
            ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
            ('FONTSIZE',(0,0),(-1,0),10),
            ('GRID',(0,0),(-1,-1),1,colors.grey),
            ('FONTSIZE',(0,1),(-1,-1),9),
            ('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white, colors.HexColor('#f3f3f3')])
        ]))
        story.append(risks_table)
    story.append(Spacer(1, 16))
    
    # Section 5: Executive Summary
    story.append(Paragraph('Executive Summary', s['Heading2']))
    story.append(Spacer(1, 8))
    bullet_style = ParagraphStyle(
        'BulletStyle',
        parent=s['Normal'],
        fontSize=10,
        leading=14,
        leftIndent=20,
        spaceAfter=6
    )
    for bullet in report.get('bullets', []):
        story.append(Paragraph(f'• {bullet}', bullet_style))
    
    doc.build(story)
    return f


def format_ppt_table(table, header_color=(46, 125, 50)):
    """Format PPT table with green header"""
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*header_color)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True
                run.font.size = Pt(11)


def export_ppt(report):
    """Export comprehensive 5-section report to PowerPoint"""
    f=os.path.join(folder(),'WSR_Report.pptx')
    prs=Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    sl=prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    week_range = report.get('week_range', 'Current Week')
    title_box = sl.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = f'Weekly Status Report : {week_range}'
    title_frame.paragraphs[0].font.size = Pt(20)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(46, 125, 50)
    
    y = 0.7
    
    # Section 1: Project Status
    if not report['summary'].empty:
        summary_title = sl.shapes.add_textbox(Inches(0.3), Inches(y), Inches(9), Inches(0.25))
        summary_title.text_frame.text = 'Project Status'
        summary_title.text_frame.paragraphs[0].font.bold = True
        summary_title.text_frame.paragraphs[0].font.size = Pt(12)
        y += 0.3
        
        sm = report['summary']
        rows, cols = len(sm) + 1, len(sm.columns)
        left, top, width, height = Inches(0.3), Inches(y), Inches(9.4), Inches(0.8)
        table = sl.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header
        for c, col in enumerate(sm.columns):
            table.cell(0, c).text = str(col)
        
        # Data
        for r in range(1, rows):
            for c in range(cols):
                table.cell(r, c).text = str(sm.iloc[r-1, c])
        
        format_ppt_table(table)
        y += 1.0
    
    # Section 2: Project Aging
    if not report['aging'].empty and y < 7:
        aging_title = sl.shapes.add_textbox(Inches(0.3), Inches(y), Inches(9), Inches(0.25))
        aging_title.text_frame.text = 'Project Aging Summary'
        aging_title.text_frame.paragraphs[0].font.bold = True
        aging_title.text_frame.paragraphs[0].font.size = Pt(12)
        y += 0.3
        
        ag = report['aging']
        rows, cols = len(ag) + 1, len(ag.columns)
        left, top, width, height = Inches(0.3), Inches(y), Inches(9.4), Inches(min(1.5, 7-y-0.3))
        table = sl.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header
        for c, col in enumerate(ag.columns):
            table.cell(0, c).text = str(col)
        
        # Data
        for r in range(1, rows):
            for c in range(cols):
                table.cell(r, c).text = str(ag.iloc[r-1, c])
        
        format_ppt_table(table)
    
    prs.save(f)
    return f


def export_png(report):
    """Export comprehensive report to PNG image"""
    f=os.path.join(folder(),'WSR_Report.png')
    img=Image.new('RGB',(1800,2200),'white')
    d=ImageDraw.Draw(img)
    
    try:
        font_title = ImageFont.truetype("arial.ttf", 28)
        font_heading = ImageFont.truetype("arial.ttf", 14)
        font_normal = ImageFont.truetype("arial.ttf", 11)
    except:
        font_title = font_heading = font_normal = ImageFont.load_default()
    
    y = 30
    week_range = report.get('week_range', 'Current Week')
    d.text((30, y), f'Weekly Status Report : {week_range}', fill=(46, 125, 50), font=font_title)
    y += 50
    
    # Executive Summary Bullets
    d.text((30, y), 'Executive Summary', fill=(46, 125, 50), font=font_heading)
    y += 35
    for bullet in report.get('bullets', []):
        d.text((50, y), f'• {bullet[:80]}...', fill='black', font=font_normal)
        y += 30
    
    y += 20
    
    # Project Status Table
    d.text((30, y), 'Project Status', fill=(46, 125, 50), font=font_heading)
    y += 35
    if not report['summary'].empty:
        x = 30
        for col in report['summary'].columns:
            d.text((x, y), str(col)[:15], fill='white', font=font_normal)
            x += 200
        y += 25
        
        x = 30
        for val in report['summary'].iloc[0].tolist():
            d.text((x, y), str(val)[:15], fill='black', font=font_normal)
            x += 200
        y += 30
    
    # Top Risk Projects
    if not report['risks'].empty:
        d.text((30, y), 'Top Risk Projects', fill=(46, 125, 50), font=font_heading)
        y += 35
        cols = list(report['risks'].columns)
        x = 30
        for col in cols:
            d.text((x, y), str(col)[:20], fill='white', font=font_normal)
            x += 300
        y += 25
        
        for idx, row in report['risks'].iterrows():
            x = 30
            for val in row.tolist():
                d.text((x, y), str(val)[:20], fill='black', font=font_normal)
                x += 300
            y += 25
    
    img.save(f)
    return f